import os
import re
import requests
import json
from pptx.enum.text import MSO_AUTO_SIZE
import re
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_DATA_LABEL_POSITION
from pptx.chart.data import CategoryChartData
from dotenv import load_dotenv
from langchain_groq import ChatGroq
from PIL import Image, ImageDraw
from io import BytesIO
from bs4 import BeautifulSoup
from datetime import datetime

# === Load environment variables ===
load_dotenv()
GROQ_API_KEY = os.getenv("GROQ_API_KEY")
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY_1")
GOOGLE_CSE_ID = os.getenv("GOOGLE_CSE_ID_1")
UNSPLASH_ACCESS_KEY = os.getenv("UNSPLASH_ACCESS_KEY")

# === Initialize LLM ===
llm = ChatGroq(
    
    groq_api_key=GROQ_API_KEY,
    model_name="llama-3.3-70b-versatile",
    
    
    # tools=[{"type":"browser_search"}]
)

# === Enhanced McKinsey Style Constants ===
MCKINSEY_COLORS = {
    "blue": RGBColor(12, 74, 126),
    "light_blue": RGBColor(79, 129, 189),
    "dark_blue": RGBColor(6, 45, 85),
    "accent_blue": RGBColor(102, 170, 238),
    "gray": RGBColor(89, 89, 89),
    "light_gray": RGBColor(217, 217, 217),
    "dark_gray": RGBColor(64, 64, 64),
    "background": RGBColor(255, 255, 255),
    "text": RGBColor(0, 0, 0),
    "white": RGBColor(255, 255, 255)
}
FONT_NAME = "Aptos Display"

# def get_slide_content_with_charts(topic, n_slides):
#     prompt = f"""
#     Create a professional PowerPoint presentation outline with exactly {n_slides} slides on the topic \"{topic}\".
#     Adopt a clear, data-driven, McKinsey-style approach.Do not include any bold text.Just plain data without decoration.

#     Begin with an Introduction:
#     Introduction: <1 paragraph overview of the topic>

#     For each slide, provide:
#     1. A "Title" (max 8 words).
#     2. A "Key Insight" (1 sentence).
#     3. A short "Context" paragraph (for slides with charts).
#     4. EITHER "Bullets" (4–5 medium detail points) OR a "Chart".

#     When using a chart, choose from:
#     - BAR, LINE, PIE, DOUGHNUT, AREA, SCATTER, STACKED_BAR, COLUMN
#     - Format chart data as JSON:
#       ```json
#       {{
#         "type": "BAR",
#         "data": [["Label", 123], ["Label 2", 456]],
#         "source": "Source: Your source here"
#       }}
#       ```

#     Format each slide like:
#     ---
#     Slide 1
#     Title: Example Title
#     Key Insight: Something meaningful
#     Context: Short description for charts
#     Chart:
#     ```json
#     {{ ... }}
#     ```
#     ---
#     OR
#     ---
#     Slide 2
#     Title: Example Title
#     Key Insight: Another insight
#     Bullets:
#     - Point 1
#     - Point 2
#     ---
#     """
#     return llm.invoke(prompt).content
def get_slide_content_with_charts(topic, n_slides):
    prompt = f"""
    You are preparing a professional, data-driven PowerPoint presentation outline with exactly {n_slides} slides on the topic "{topic}".
    Use a clear, McKinsey-style approach. No bold, italics, underlines, Markdown syntax, asterisks, or any other decorative formatting — only plain text and data.
    Do not use bold, italics, underlines, Markdown symbols (like *, **, _), or any decorative formatting.
      
    **Recency requirement:**
    - Use the browser search tool to find and include the most recent (up to the current month/year) credible information from authoritative sources.
    - Incorporate up-to-date statistics, trends, market data, and developments relevant to the topic.
    - Always include a source and year for any data or factual claim.

    Begin with:
    Introduction: <1-paragraph overview with latest trends and facts>

    For each slide, provide:
    1. "Title" (max 8 words).
    2. "Key Insight" (1 sentence, backed by recent facts).
    3. For chart slides: a "Context" paragraph explaining why this recent data matters.
    4. EITHER "Bullets" (4–5 medium-detail, current points) OR a "Chart".

    Chart guidelines:
    - Allowed types: BAR, LINE, PIE, DOUGHNUT, AREA, SCATTER, STACKED_BAR, COLUMN.
    - Format chart data as JSON:
      ```json
      {{
        "type": "BAR",
        "data": [["Label", 123], ["Label 2", 456]],
        "source": "Source: Organization / publication name, YYYY"
      }}
      ```

    Format each slide exactly like:
    ---
    Slide 1
    Title: Example Title
    Key Insight: Something meaningful and recent
    Context: Short description for charts
    Chart:
    ```json
    {{ ... }}
    ```
    ---
    OR
    ---
    Slide 2
    Title: Example Title
    Key Insight: Another insight
    Bullets:
    - Point 1
    - Point 2
    ---
    """
    return llm.invoke(prompt).content


def parse_mckinsey_response(text):
    slides = []
    intro_match = re.search(r"Introduction:\s*(.*?)\n---", text, re.DOTALL)
    introduction = ""

    slide_chunks = text.split('---')
    for chunk in slide_chunks:
        if not chunk.strip():
            continue

        try:
            title_match = re.search(r"Title: (.*?)\n", chunk, re.IGNORECASE)
            insight_match = re.search(r"Key Insight: (.*?)\n", chunk, re.IGNORECASE)
            context_match = re.search(r"Context: (.*?)\n", chunk, re.IGNORECASE)

            title = title_match.group(1).strip() if title_match else f"Slide {len(slides)+1}"
            insight = insight_match.group(1).strip() if insight_match else None
            context = context_match.group(1).strip() if context_match else ""

            if not insight:
                print(f"Skipping slide due to missing insight: {chunk[:100]}")
                continue

            content = {"title": title, "insight": insight, "context": context}

            chart_match = re.search(r"Chart:\s*```json\n(.*?)\n```", chunk, re.DOTALL)
            if chart_match:
                try:
                    chart_data = json.loads(chart_match.group(1).strip())
                    content["type"] = "chart"
                    content["data"] = chart_data
                except json.JSONDecodeError:
                    print(f"[Warning] Chart JSON malformed. Falling back to bullets.")
                    content["type"] = "bullets"
                    content["data"] = ["No chart data was available."]
            else:
                bullets_match = re.search(r"Bullets:\s*(.*)", chunk, re.DOTALL)
                bullets_raw = bullets_match.group(1) if bullets_match else ""
                bullets = [b.strip("- ").strip() for b in bullets_raw.split('\n') if b.strip()]
                if not bullets:
                    bullets = ["No content was provided by the AI for this slide."]
                content["type"] = "bullets"
                content["data"] = bullets

            slides.append(content)
        except Exception as e:
            print(f"Could not parse a slide chunk, skipping. Error: {e}\nChunk: {chunk[:100]}...")
            continue

    return {"intro": introduction, "slides": slides}

def add_enhanced_title_slide(prs, topic, intro):
    """Creates a stunning McKinsey-style title slide with professional layout and design elements."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = MCKINSEY_COLORS["background"]

    # Background geometric elements for visual appeal
    bg_rect = slide.shapes.add_shape(1, Inches(8), Inches(0), Inches(3.33), Inches(7.5))
    bg_rect.fill.solid()
    bg_rect.fill.fore_color.rgb = MCKINSEY_COLORS["light_blue"]
    bg_rect.line.fill.background()
    
    # Diagonal accent shape
    accent_shape = slide.shapes.add_shape(1, Inches(0), Inches(6.7), Inches(13.33), Inches(1.0))
    accent_shape.fill.solid()
    accent_shape.fill.fore_color.rgb = MCKINSEY_COLORS["blue"]
    accent_shape.line.fill.background()

    # Main Title - Bold and prominent
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(10), Inches(2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = topic.upper()
    p.font.name = FONT_NAME
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = MCKINSEY_COLORS["dark_blue"]
    p.alignment = PP_ALIGN.LEFT

    # Professional subtitle with current date
    current_date = datetime.now().strftime("%B %Y")
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(10), Inches(0.8))
    tf2 = subtitle_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = f"Strategic Analysis & Business Intelligence | {current_date}"
    p2.font.name = FONT_NAME
    p2.font.size = Pt(20)
    p2.font.color.rgb = MCKINSEY_COLORS["gray"]
    p2.alignment = PP_ALIGN.LEFT

    # Executive Summary Box with better formatting
    intro_box = slide.shapes.add_textbox(Inches(1), Inches(4.0), Inches(10), Inches(1.5))
    tf3 = intro_box.text_frame
    tf3.word_wrap = True
    # tf3.auto_size = False
    tf3.auto_size = MSO_AUTO_SIZE.NONE
    tf3.margin_left = Inches(0.2)
    tf3.margin_right = Inches(0.2)
    tf3.margin_top = Inches(0.1)
    tf3.margin_bottom = Inches(0.1)
    
    p3 = tf3.paragraphs[0]
    p3.font.name = FONT_NAME
    p3.font.size = Pt(14)
    p3.font.bold = True
    p3.font.color.rgb = MCKINSEY_COLORS["blue"]
    p3.alignment = PP_ALIGN.LEFT
    
    # Add the actual introduction content
    p4 = tf3.add_paragraph()
    p4.text = intro
    p4.font.name = FONT_NAME
    p4.font.size = Pt(16)
    p4.font.color.rgb = MCKINSEY_COLORS["text"]
    p4.alignment = PP_ALIGN.LEFT
    p4.space_before = Pt(8)

    # Professional branding footer
    brand_box = slide.shapes.add_textbox(Inches(1), Inches(7), Inches(6), Inches(0.4))
    tf4 = brand_box.text_frame
    p5 = tf4.paragraphs[0]
    p5.text = "MADE BY ENTHRAL AI"
    p5.font.name = FONT_NAME
    p5.font.size = Pt(10)
    p5.font.italic = True
    p5.font.color.rgb = MCKINSEY_COLORS["white"]
    p5.alignment = PP_ALIGN.LEFT

    # Decorative elements
    line1 = slide.shapes.add_shape(1, Inches(1), Inches(4.2), Inches(4), Inches(0))
    line1.line.color.rgb = MCKINSEY_COLORS["accent_blue"]
    line1.line.width = Pt(4)
    
    line2 = slide.shapes.add_shape(1, Inches(1), Inches(4.3), Inches(2), Inches(0))
    line2.line.color.rgb = MCKINSEY_COLORS["blue"]
    line2.line.width = Pt(2)

def add_chart_slide_with_context(slide, chart_info, context):
    if context:
        tx_box = slide.shapes.add_textbox(Inches(1), Inches(1.6), Inches(11), Inches(0.6))
        tf = tx_box.text_frame
        p = tf.paragraphs[0]
        p.text = context
        p.font.size = Pt(14)
        p.font.name = FONT_NAME
        p.font.color.rgb = MCKINSEY_COLORS["gray"]

    add_chart_slide(slide, chart_info)

def calculate_dynamic_image_size(image_path, max_width=4.5, max_height=4.0):
    """Calculate optimal image size while maintaining aspect ratio and fitting within constraints."""
    try:
        with Image.open(image_path) as img:
            original_width, original_height = img.size
            aspect_ratio = original_width / original_height
            
            # Calculate dimensions based on aspect ratio
            if aspect_ratio > (max_width / max_height):
                # Image is wider - constrain by width
                width = max_width
                height = max_width / aspect_ratio
            else:
                # Image is taller - constrain by height
                height = max_height
                width = max_height * aspect_ratio
            
            # Ensure minimum size for visibility
            width = max(width, 2.0)
            height = max(height, 1.5)
            
            return Inches(width), Inches(height)
    except Exception as e:
        print(f"Error calculating image size: {e}")
        return Inches(3.5), Inches(2.8)  # Default fallback size

def build_mckinsey_ppt(parsed_data, topic):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    add_enhanced_title_slide(prs, topic, parsed_data["intro"])

    for i, slide_content in enumerate(parsed_data["slides"], start=1):
        slide = prs.slides.add_slide(blank_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = MCKINSEY_COLORS["background"]

        add_enhanced_header(slide, slide_content['title'], slide_content['insight'])
        add_enhanced_footer(slide, i + 1)

        if slide_content['type'] == 'chart':
            add_chart_slide_with_context(slide, slide_content['data'], slide_content.get("context", ""))
        else:
            image_path = fetch_image(f"{slide_content['title']} {slide_content['insight']}", topic)
            add_enhanced_text_and_image_slide(slide, slide_content['data'], image_path)

    filename = topic.strip().replace(" ", "_") + "_McKinsey_Style.pptx"
    prs.save(filename)
    print(f"\nPresentation saved as: {filename}")

def add_enhanced_header(slide, title, insight):
    """Adds a professionally styled title and key insight with better visual hierarchy."""
    # Background header area
    header_bg = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.4))
    header_bg.fill.solid()
    header_bg.fill.fore_color.rgb = MCKINSEY_COLORS["background"]
    header_bg.line.color.rgb = MCKINSEY_COLORS["light_gray"]
    header_bg.line.width = Pt(1)
    
    # Title
    tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.65))
    tf = tx_box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = title.upper()
    p.font.name = FONT_NAME
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = MCKINSEY_COLORS["dark_blue"]

    # Key Insight with better styling
    insight_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(12), Inches(0.5))
    tf2 = insight_box.text_frame
    tf2.margin_left = Inches(0.1)
    p_insight = tf2.paragraphs[0]
    p_insight.text = f"Key Insight: {insight}"
    p_insight.font.name = FONT_NAME
    p_insight.font.size = Pt(14)
    p_insight.font.color.rgb = MCKINSEY_COLORS["gray"]
    p_insight.font.italic = True

    # Enhanced blue accent line
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.35), Inches(12), Inches(0))
    line.line.color.rgb = MCKINSEY_COLORS["blue"]
    line.line.width = Pt(3)

def add_enhanced_footer(slide, slide_number):
    """Adds a professionally styled footer with slide number and branding."""
    # Footer background
    footer_bg = slide.shapes.add_shape(1, Inches(0), Inches(6.9), Inches(13.33), Inches(0.6))
    footer_bg.fill.solid()
    footer_bg.fill.fore_color.rgb = MCKINSEY_COLORS["light_gray"]
    footer_bg.line.fill.background()
    
    # Slide number
    tx_box = slide.shapes.add_textbox(Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.4))
    tf = tx_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{slide_number}"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FONT_NAME
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = MCKINSEY_COLORS["blue"]
    
    # Professional footer text
    footer_text_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(6), Inches(0.4))
    tf2 = footer_text_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "MADE BY ENTHRAL AI"
    p2.font.name = FONT_NAME
    p2.font.size = Pt(9)
    p2.font.color.rgb = MCKINSEY_COLORS["dark_gray"]

# def fetch_image(prompt: str, topic: str) -> str:
#     from urllib.parse import quote
#     headers = {"User-Agent": "Mozilla/5.0"}
#     query = f"{topic} {prompt}".strip()
#     encoded = quote(query)

#     def is_vertical(img):
#         w, h = img.size
#         aspect_ratio = w / h
#         return aspect_ratio < 1.25

#     def try_save_image(url, label, require_vertical=True):
#         try:
#             r = requests.get(url, stream=True, headers=headers, timeout=10)
#             if r.status_code == 200 and 'image' in r.headers.get('Content-Type', ''):
#                 img = Image.open(BytesIO(r.content)).convert("RGB")
#                 if not require_vertical or is_vertical(img):
#                     filename = f"{label}_{re.sub(r'[^a-zA-Z0-9]', '_', url)[:30]}.jpg"
#                     img.save(filename, format="JPEG")
#                     return filename
#         except:
#             pass
#         return None

#     def fetch_from_google(require_vertical=True):
#         try:
#             g_url = f"https://www.googleapis.com/customsearch/v1?q={encoded}&searchType=image&num=5&key={GOOGLE_API_KEY}&cx={GOOGLE_CSE_ID}"
#             response = requests.get(g_url, timeout=10)
#             data = response.json()
#             if 'items' in data:
#                 for item in data['items']:
#                     url = item.get("link", "")
#                     if url:
#                         img_path = try_save_image(url, "google", require_vertical)
#                         if img_path:
#                             return img_path
#         except:
#             pass
#         return None

#     for fetch_fn in [fetch_from_google]:
#         image_path = fetch_fn(require_vertical=True)
#         if image_path:
#             return image_path
#         image_path = fetch_fn(require_vertical=False)
#         if image_path:
#             return image_path

#     return "fallback.jpg"

# def validated_image_bytes(path, max_width=800, max_height=500):
#     try:
#         with Image.open(path) as img:
#             img.verify()
#         with Image.open(path) as img:
#             img = img.convert("RGB")
#             img.thumbnail((max_width, max_height))
#             stream = BytesIO()
#             img.save(stream, format="PNG")
#             stream.seek(0)
#             return stream
#     except Exception as e:
#         print(f"[Image Validation Error] {e} - Path: {path}")
#         return None

def fetch_image(prompt: str, topic: str) -> str:
    """
    Enhanced image fetching with multiple sources:
    1. Google Custom Search API
    2. Unsplash API (fallback)
    3. fallback.jpg (final fallback)
    """
    from urllib.parse import quote
    import time
    
    headers = {"User-Agent": "Mozilla/5.0"}
    query = f"{topic} {prompt}".strip()
    encoded = quote(query)

    def is_vertical(img):
        w, h = img.size
        aspect_ratio = w / h
        return aspect_ratio < 1.25

    def try_save_image(url, label, require_vertical=True):
        try:
            r = requests.get(url, stream=True, headers=headers, timeout=10)
            if r.status_code == 200 and 'image' in r.headers.get('Content-Type', ''):
                img = Image.open(BytesIO(r.content)).convert("RGB")
                if not require_vertical or is_vertical(img):
                    filename = f"{label}_{re.sub(r'[^a-zA-Z0-9]', '_', url)[:30]}.jpg"
                    img.save(filename, format="JPEG")
                    print(f"✓ Image saved from {label}: {filename}")
                    return filename
        except Exception as e:
            print(f"✗ Failed to save image from {label}: {e}")
        return None

    def fetch_from_google(require_vertical=True):
        """Fetch images from Google Custom Search API"""
        try:
            if not GOOGLE_API_KEY or not GOOGLE_CSE_ID:
                print("✗ Google API credentials not configured")
                return None
                
            g_url = f"https://www.googleapis.com/customsearch/v1?q={encoded}&searchType=image&num=8&key={GOOGLE_API_KEY}&cx={GOOGLE_CSE_ID}"
            response = requests.get(g_url, timeout=15)
            
            if response.status_code != 200:
                print(f"✗ Google API error: {response.status_code}")
                return None
                
            data = response.json()
            
            if 'items' not in data:
                print("✗ No Google search results found")
                return None
                
            print(f"→ Found {len(data['items'])} Google images, trying to download...")
            
            for i, item in enumerate(data['items']):
                url = item.get("link", "")
                if url:
                    img_path = try_save_image(url, f"google_{i}", require_vertical)
                    if img_path:
                        return img_path
                        
        except Exception as e:
            print(f"✗ Google search failed: {e}")
        return None

    def fetch_from_unsplash(require_vertical=True):
        """Fetch images from Unsplash API"""
        try:
            if not UNSPLASH_ACCESS_KEY:
                print("✗ Unsplash API key not configured")
                return None
            
            # Clean and optimize search query for Unsplash
            unsplash_query = query.replace("slide", "").replace("presentation", "").strip()
            unsplash_encoded = quote(unsplash_query)
            
            # Unsplash API endpoint
            unsplash_url = f"https://api.unsplash.com/search/photos"
            unsplash_headers = {
                "Authorization": f"Client-ID {UNSPLASH_ACCESS_KEY}",
                "User-Agent": "Mozilla/5.0 (PowerPoint Generator)"
            }
            
            params = {
                "query": unsplash_query,
                "per_page": 10,
                "orientation": "portrait" if require_vertical else "landscape",
                "order_by": "relevant"
            }
            
            response = requests.get(unsplash_url, headers=unsplash_headers, params=params, timeout=15)
            
            if response.status_code != 200:
                print(f"✗ Unsplash API error: {response.status_code}")
                return None
            
            data = response.json()
            results = data.get("results", [])
            
            if not results:
                print("✗ No Unsplash results found")
                return None
                
            print(f"→ Found {len(results)} Unsplash images, trying to download...")
            
            for i, photo in enumerate(results):
                try:
                    # Get the regular size image URL
                    img_url = photo["urls"]["regular"]  # Good balance of quality and size
                    
                    # Try to download and save
                    img_path = try_save_image(img_url, f"unsplash_{i}", require_vertical)
                    if img_path:
                        return img_path
                        
                except KeyError as e:
                    print(f"✗ Unsplash photo data missing key: {e}")
                    continue
                    
        except Exception as e:
            print(f"✗ Unsplash search failed: {e}")
        return None

    def create_better_fallback():
        """Create an enhanced fallback image if none exists"""
        try:
            # Create a more professional fallback image
            img = Image.new('RGB', (800, 600), color=(245, 245, 245))  # Light gray background
            draw = ImageDraw.Draw(img)
            
            # Add McKinsey-style design elements
            # Main border
            draw.rectangle([40, 40, 760, 560], outline=(12, 74, 126), width=4)  # McKinsey blue
            
            # Inner design elements
            draw.rectangle([80, 80, 720, 520], outline=(217, 217, 217), width=2)  # Light gray
            
            # Add text
            try:
                # Try to load a decent font, fallback to default if not available
                from PIL import ImageFont
                try:
                    font_large = ImageFont.truetype("arial.ttf", 36)
                    font_small = ImageFont.truetype("arial.ttf", 18)
                except:
                    font_large = ImageFont.load_default()
                    font_small = ImageFont.load_default()
            except:
                font_large = ImageFont.load_default()
                font_small = ImageFont.load_default()
            
            # Main text
            draw.text((400, 250), "Professional", fill=(12, 74, 126), anchor='mm', font=font_large)
            draw.text((400, 300), "Image Placeholder", fill=(12, 74, 126), anchor='mm', font=font_large)
            
            # Subtitle
            draw.text((400, 380), f"Topic: {topic}", fill=(89, 89, 89), anchor='mm', font=font_small)
            
            # Add some geometric elements for visual appeal
            draw.ellipse([320, 150, 480, 200], outline=(79, 129, 189), width=3)
            
            img.save('fallback.jpg', 'JPEG', quality=90)
            print("✓ Created enhanced fallback image")
            return True
            
        except Exception as e:
            print(f"✗ Failed to create enhanced fallback: {e}")
            return False

    # Main fetching logic with cascading fallbacks
    print(f"\n Searching for images: '{query}'")
    
    # Method 1: Google Custom Search (try both vertical and any orientation)
    for orientation_pref in [True, False]:  # Try vertical first, then any
        print(f"\n→ Trying Google Custom Search ({'vertical preferred' if orientation_pref else 'any orientation'})...")
        image_path = fetch_from_google(require_vertical=orientation_pref)
        if image_path:
            return image_path
    
    # Method 2: Unsplash API (try both orientations)
    for orientation_pref in [True, False]:  # Try vertical first, then any
        print(f"\n→ Trying Unsplash API ({'vertical preferred' if orientation_pref else 'any orientation'})...")
        image_path = fetch_from_unsplash(require_vertical=orientation_pref)
        if image_path:
            return image_path
    
    # Method 3: Final fallback
    print(f"\n→ Using fallback image...")
    
    # Check if fallback exists, if not create a better one
    if not os.path.exists("fallback.jpg"):
        print("→ Creating enhanced fallback image...")
        create_better_fallback()
    
    if os.path.exists("fallback.jpg"):
        print("✓ Using fallback.jpg")
        return "fallback.jpg"
    else:
        # Create a very basic fallback if all else fails
        print("→ Creating basic emergency fallback...")
        try:
            img = Image.new('RGB', (800, 600), color='lightgray')
            draw = ImageDraw.Draw(img)
            draw.rectangle([50, 50, 750, 550], outline='darkgray', width=3)
            draw.text((400, 300), "Image Not Available", fill='darkgray', anchor='mm')
            img.save('fallback.jpg')
            return "fallback.jpg"
        except Exception as e:
            print(f"✗ Emergency fallback creation failed: {e}")
            return None


def validated_image_bytes(path, max_width=800, max_height=500):
    """
    Enhanced image validation with better error handling and optimization.
    """
    try:
        if not path or not os.path.exists(path):
            print(f"✗ Image path does not exist: {path}")
            return None
            
        # First verify the image can be opened
        with Image.open(path) as img:
            img.verify()
            
        # Now process the image
        with Image.open(path) as img:
            # Convert to RGB if necessary (handles RGBA, P, etc.)
            if img.mode != 'RGB':
                img = img.convert("RGB")
            
            # Get original dimensions
            original_w, original_h = img.size
            print(f"→ Processing image: {original_w}x{original_h} -> ", end="")
            
            # Resize if needed while maintaining aspect ratio
            img.thumbnail((max_width, max_height), Image.Resampling.LANCZOS)
            final_w, final_h = img.size
            print(f"{final_w}x{final_h}")
            
            # Convert to bytes
            stream = BytesIO()
            img.save(stream, format="PNG", optimize=True, quality=90)
            stream.seek(0)
            
            return stream
            
    except Exception as e:
        print(f"✗ Image validation error for {path}: {e}")
        return None

def add_chart_slide(slide, chart_info):
    """Adds a dynamically generated chart to the slide, supporting both single and multi-series data."""
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import (
        XL_CHART_TYPE, XL_LEGEND_POSITION, XL_DATA_LABEL_POSITION
    )

    try:
        raw_data = chart_info.get("data", [])
        if not raw_data or not isinstance(raw_data, list):
            raise ValueError("Chart data is missing or malformed.")

        chart_type_str = chart_info.get("type", "BAR").upper()
        chart_type_map = {
            "BAR": XL_CHART_TYPE.BAR_CLUSTERED,
            "COLUMN": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "LINE": XL_CHART_TYPE.LINE,
            "PIE": XL_CHART_TYPE.PIE,
            "DOUGHNUT": XL_CHART_TYPE.DOUGHNUT,
            "AREA": XL_CHART_TYPE.AREA,
            "SCATTER": XL_CHART_TYPE.XY_SCATTER_LINES,
            "STACKED_BAR": XL_CHART_TYPE.BAR_STACKED
        }
        chart_type = chart_type_map.get(chart_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED)

        chart_data = CategoryChartData()

        # Multi-Series (e.g., ['Region', 40, 30])
        if all(isinstance(row, list) and len(row) > 2 for row in raw_data):
            categories = [row[0] for row in raw_data]
            num_series = len(raw_data[0]) - 1
            series_data = [[] for _ in range(num_series)]

            for row in raw_data:
                for i in range(num_series):
                    val = row[i + 1]
                    series_data[i].append(float(val) if isinstance(val, (int, float)) else 0)

            chart_data.categories = categories
            for i, values in enumerate(series_data):
                chart_data.add_series(f"Series {i + 1}", values)

        # Single-Series (e.g., ['Label', value])
        elif all(isinstance(row, list) and len(row) == 2 for row in raw_data):
            chart_data.categories = [str(row[0]) for row in raw_data]
            values = [float(row[1]) if isinstance(row[1], (int, float)) else 0.0 for row in raw_data]
            chart_data.add_series("", values)

        else:
            raise ValueError("Each chart data point must be a [label, value] pair or valid multi-series row.")

        # Draw chart
        x, y, cx, cy = Inches(1), Inches(1.7), Inches(11), Inches(4.5)
        graphic_frame = slide.shapes.add_chart(chart_type, x, y, cx, cy, chart_data)
        chart = graphic_frame.chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False

        plot = chart.plots[0]
        plot.has_data_labels = True
        plot.data_labels.font.size = Pt(12)
        plot.data_labels.font.bold = True
        plot.data_labels.font.color.rgb = MCKINSEY_COLORS["gray"]

        # Axis styling
        axis_supported_types = {
            XL_CHART_TYPE.BAR_CLUSTERED,
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            XL_CHART_TYPE.BAR_STACKED,
            XL_CHART_TYPE.LINE,
            XL_CHART_TYPE.AREA
        }
        if chart.chart_type in axis_supported_types:
            try:
                chart.category_axis.tick_labels.font.size = Pt(11)
                chart.value_axis.has_major_gridlines = False
                chart.value_axis.tick_labels.font.size = Pt(11)
                plot.data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
            except AttributeError:
                pass
        elif chart_type in (XL_CHART_TYPE.PIE, XL_CHART_TYPE.DOUGHNUT):
            plot.data_labels.position = XL_DATA_LABEL_POSITION.BEST_FIT

        # Chart source
        source_text = chart_info.get("source", "")
        if source_text:
            tx_box = slide.shapes.add_textbox(Inches(1), Inches(6.4), Inches(11), Inches(0.4))
            p = tx_box.text_frame.paragraphs[0]
            p.text = source_text
            p.font.size = Pt(9)
            p.font.italic = True
            p.font.color.rgb = MCKINSEY_COLORS["gray"]

    except Exception as e:
        print(f"Failed to create chart: {e}. Data: {chart_info}")
        tx_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(1))
        tx_box.text_frame.paragraphs[0].text = "Error: Could not generate the requested chart."

# def add_enhanced_text_and_image_slide(slide, bullets, image_path):
#     """Adds professionally formatted bullet points and properly positioned image with better layout."""
    
#     # FIXED: Better layout with proper spacing between text and image
#     # Text Box (Left side) - Reduced width to provide more space
#     tx_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.7), Inches(5.8), Inches(4.8))
#     tf = tx_box.text_frame
#     tf.word_wrap = True
#     tf.margin_left = Inches(0.2)
#     tf.margin_right = Inches(0.1)
#     tf.margin_top = Inches(0.1)
#     tf.margin_bottom = Inches(0.1)

#     for i, bullet_text in enumerate(bullets):
#         p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        
#         # Add bullet symbol with better formatting
#         p.text = f"• {bullet_text}"
#         p.font.name = FONT_NAME
#         p.font.size = Pt(16)  # Slightly smaller to fit better
#         p.font.color.rgb = MCKINSEY_COLORS["text"]
#         p.level = 0
#         p.space_before = Pt(12)
#         p.space_after = Pt(6)
#         p.line_spacing = 1.15
        
#         # Make bullet blue for visual appeal
#         if len(p.text) > 2:
#             p.runs[0].font.color.rgb = MCKINSEY_COLORS["blue"]
#             p.runs[0].font.bold = True

#     # FIXED: Image positioning with proper spacing and centering
#     img_stream = validated_image_bytes(image_path)
#     if not img_stream:
#         img_stream = validated_image_bytes("fallback.jpg")

#     if img_stream:
#         try:
#             # Calculate dynamic size with better constraints
#             width, height = calculate_dynamic_image_size(image_path, max_width=4.2, max_height=3.8)
            
#             # FIXED: Better positioning with proper spacing from text
#             # Image starts at x=7.2 (giving 0.6 inches gap from text ending at 6.6)
#             x_position = Inches(7.2)
            
#             # Center the image vertically in the available content space
#             available_height = Inches(4.8)
#             content_start = Inches(1.7)
#             y_position = content_start + (available_height - height) / 2
            
#             # Add image with calculated dimensions
#             pic = slide.shapes.add_picture(
#                 img_stream, 
#                 x_position, 
#                 y_position, 
#                 width=width, 
#                 height=height
#             )
            
#             # Add professional border around image
#             pic.line.color.rgb = MCKINSEY_COLORS["light_gray"]
#             pic.line.width = Pt(2)
            
#             # FIXED: Add subtle shadow effect by creating a background rectangle
#             shadow_rect = slide.shapes.add_shape(
#                 1,  # Rectangle shape
#                 x_position + Inches(0.05),  # Slightly offset
#                 y_position + Inches(0.05),  # Slightly offset
#                 width,
#                 height
#             )
#             shadow_rect.fill.solid()
#             shadow_rect.fill.fore_color.rgb = MCKINSEY_COLORS["light_gray"]
#             shadow_rect.line.fill.background()
            
#             # Move shadow behind the image
#             shadow_rect.element.getparent().remove(shadow_rect.element)
#             slide._element.insert(-1, shadow_rect.element)
            
#         except Exception as e:
#             print(f"[Image Insert Error] {e} - Using fallback positioning")
#             # Fallback with basic positioning
#             try:
#                 pic = slide.shapes.add_picture(
#                     img_stream if img_stream else validated_image_bytes("fallback.jpg"), 
#                     Inches(7.2), 
#                     Inches(2.5), 
#                     width=Inches(4.0), 
#                     height=Inches(3.0)
#                 )
#                 pic.line.color.rgb = MCKINSEY_COLORS["light_gray"]
#                 pic.line.width = Pt(1)
#             except Exception as e2:
#                 print(f"[Fallback Image Error] {e2}")
def classify_paragraph_type(raw_text):
    """
    Classifies the type of presentation based on user input text.
    Categories: business, academic, technical, educational, motivational, general.
    """
    prompt = f"""
    Classify the following text into one category:
    - business
    - academic
    - technical
    - educational
    - motivational
    - general

    Text:
    \"\"\"{raw_text}\"\"\"

    Respond with only one word (the category).
    """
    result = llm.invoke(prompt).content.strip().lower()
    if result not in ["business", "academic", "technical", "educational", "motivational", "general"]:
        return "business"  # default fallback
    return result
def refine_paragraph_input(raw_text, category="business"):
    """
    Refines raw text differently depending on content type.
    """
    style_guidelines = {
        "business": "Make it concise, factual, and strategic. Focus on insights, numbers, and implications.",
        "academic": "Make it clear, explanatory, with definitions and logical flow.",
        "technical": "Make it precise, structured, with focus on processes, systems, and technical clarity.",
        "educational": "Make it simple, clear, and beginner-friendly with learning outcomes.",
        "motivational": "Make it inspiring, story-driven, with positive tone and key messages.",
        "general": "Make it well-structured, clear, and neutral."
    }

    guidelines = style_guidelines.get(category, style_guidelines["business"])

    prompt = f"""
    Refine and summarize the following text for a professional {category} presentation.
    {guidelines}
    Output 5-6 concise paragraphs, well-structured for slide generation.
    Do not use bold, italics, underlines, Markdown symbols (like *, **, _), or any decorative formatting.
    Only plain text should be returned.
    Text:
    \"\"\"{raw_text}\"\"\"

    Refined version:
    """
    return llm.invoke(prompt).content
def get_slide_content_from_paragraph(context_text, category="business", min_slides=1, max_slides=15):
    style_instructions = {
        "business": "Use a McKinsey-style with charts, insights, and data-driven points.",
        "academic": "Use an academic style with definitions, theories, and structured explanation.",
        "technical": "Use a technical style with system diagrams, architecture, workflows, or pseudocode.",
        "educational": "Use an educational style with simple language, learning outcomes, and key concepts.",
        "motivational": "Use a motivational style with quotes, storytelling, and call-to-action messages.",
        "general": "Use a general informative style with clarity and balanced explanation."
    }
    style = style_instructions.get(category, style_instructions["business"])

    prompt = f"""
    You are preparing a professional {category} PowerPoint presentation outline 
    based on the following refined content:

    \"\"\"{context_text}\"\"\"

    Rules:
    - Create between {min_slides} and {max_slides} slides depending on content richness.
    - {style}
    - Return ONLY valid JSON, no explanations or markdown.
    - "intro" must be maximum 2-3 lines.
    - Each slide must include:
      - "title": max 8 words
      - "insight": one or two sentences summarizing the key takeaway
      - "type": either "bullets" or "chart"
      - If "type" = "bullets", "data" must be an array of objects:
        {{"point": "short phrase", "desc": "5-6 line proper explanation"}}
      - If "type" = "chart", "data" must be JSON:
        {{
          "type": "BAR" | "LINE" | "PIE" | "COLUMN" | "DOUGHNUT" | "AREA" | "SCATTER" | "STACKED_BAR",
          "data": [["Label1", 123], ["Label2", 456]],
          "source": "Source: Organization, 2025"
        }}
        and optionally include "context": "Why this data matters"

    JSON format:
    {{
      "intro": "short intro text",
      "slides": [
        {{
          "title": "Slide Title",
          "insight": "Key insight",
          "type": "bullets",
          "data": [
            {{"point": "Clarity", "desc": "Use straightforward language to avoid confusion"}},
            {{"point": "Audience Awareness", "desc": "Adapt tone and style to the readers"}}
          ]
        }},
        {{
          "title": "Slide with Chart",
          "insight": "Data-driven point",
          "type": "chart",
          "data": {{
            "type": "COLUMN",
            "data": [["X", 10], ["Y", 20]],
            "source": "Source: Example Org, 2025"
          }},
          "context": "Why this data is important"
        }}
      ]
    }}
    """
    return llm.invoke(prompt).content

def parse_json_slides(text):
    try:
        return json.loads(text)
    except json.JSONDecodeError as e:
        print(f"[Error] JSON parse failed: {e}\n{text[:300]}...")
        return {"intro": "", "slides": []}


# def add_enhanced_text_and_image_slide(slide, bullets, image_path):
#     """
#     Creates an elegant side-by-side layout with bullet points on the left 
#     and a large, professionally positioned image on the right.
#     """
    
#     # Left side - Text content (takes up ~45% of slide width)
#     text_width = Inches(5.5)  # Increased width for better text space
#     text_height = Inches(4.5)
    
#     tx_box = slide.shapes.add_textbox(
#         Inches(0.8),      # Left margin
#         Inches(1.8),      # Top margin (below header)
#         text_width,       # Width
#         text_height       # Height
#     )
    
#     tf = tx_box.text_frame
#     tf.word_wrap = True
#     tf.margin_left = Inches(0.15)
#     tf.margin_right = Inches(0.15)
#     tf.margin_top = Inches(0.1)
#     tf.margin_bottom = Inches(0.1)
#     # tf.auto_size = True
#     tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

#     # Add bullet points with enhanced styling
#     for i, bullet_text in enumerate(bullets):
#         if i > 0:
#             p = tf.add_paragraph()
#         else:
#             p = tf.paragraphs[0]
        
#         # Clean bullet formatting
#         p.text = f"• {bullet_text.strip()}"
#         p.font.name = FONT_NAME
#         p.font.size = Pt(16)
#         p.font.color.rgb = MCKINSEY_COLORS["text"]
#         p.level = 0
        
#         # Spacing for better readability
#         p.space_before = Pt(8)
#         p.space_after = Pt(12)
#         p.line_spacing = 1.2
        
#         # Style the bullet point
#         if len(p.runs) > 0:
#             # Make bullet symbol blue and bold
#             bullet_run = p.runs[0]
#             bullet_run.font.color.rgb = MCKINSEY_COLORS["blue"]
#             bullet_run.font.bold = True

#     # Right side - Image (takes up ~50% of slide width)
#     img_stream = validated_image_bytes(image_path)
#     if not img_stream:
#         img_stream = validated_image_bytes("fallback.jpg")

#     if img_stream:
#         try:
#             # Calculate optimal image dimensions for the right side
#             available_width = Inches(5.8)   # Generous width for impact
#             available_height = Inches(4.5)  # Match text height
            
#             # Get actual image dimensions to maintain aspect ratio
#             with Image.open(image_path if os.path.exists(image_path) else "fallback.jpg") as img:
#                 original_width, original_height = img.size
#                 aspect_ratio = original_width / original_height
                
#                 # Calculate best fit dimensions
#                 if aspect_ratio > (available_width.inches / available_height.inches):
#                     # Image is wider - fit by width
#                     img_width = available_width
#                     img_height = Inches(available_width.inches / aspect_ratio)
#                 else:
#                     # Image is taller - fit by height  
#                     img_height = available_height
#                     img_width = Inches(available_height.inches * aspect_ratio)
                
#                 # Ensure minimum size for visibility
#                 img_width = max(img_width, Inches(4.0))
#                 img_height = max(img_height, Inches(3.0))
            
#             # Position image on the right side with proper spacing
#             img_x = Inches(7.0)  # Start after text with spacing
#             img_y = Inches(1.8) + (available_height - img_height) / 2  # Center vertically
            
#             # Add the main image
#             pic = slide.shapes.add_picture(
#                 img_stream,
#                 img_x,
#                 img_y,
#                 width=img_width,
#                 height=img_height
#             )
            
#             # Add subtle border for professional look
#             pic.line.color.rgb = MCKINSEY_COLORS["light_gray"]
#             pic.line.width = Pt(1)
            
#             # Create subtle shadow effect
#             shadow_offset = Inches(0.08)
#             shadow_rect = slide.shapes.add_shape(
#                 1,  # Rectangle shape
#                 img_x + shadow_offset,
#                 img_y + shadow_offset,
#                 img_width,
#                 img_height
#             )
#             shadow_rect.fill.solid()
#             shadow_rect.fill.fore_color.rgb = RGBColor(200, 200, 200)  # Light gray shadow
#             shadow_rect.line.fill.background()
            
#             # Move shadow behind the image (z-order)
#             shadow_rect.element.getparent().remove(shadow_rect.element)
#             slide._element.insert(-1, shadow_rect.element)
            
#         except Exception as e:
#             print(f"[Image Error] {e} - Using fallback layout")
#             # Fallback: simpler positioning
#             try:
#                 pic = slide.shapes.add_picture(
#                     img_stream,
#                     Inches(7.0),
#                     Inches(2.0),
#                     width=Inches(5.0),
#                     height=Inches(3.5)
#                 )
#                 pic.line.color.rgb = MCKINSEY_COLORS["light_gray"]
#                 pic.line.width = Pt(1)
#             except Exception as e2:
#                 print(f"[Fallback Image Error] {e2}")
    
#     # Optional: Add a subtle dividing line between text and image
#     divider_line = slide.shapes.add_shape(
#         1,  # Line shape
#         Inches(6.5),   # X position (between text and image)
#         Inches(1.8),   # Y start
#         Inches(0),     # Width (vertical line)
#         Inches(4.5)    # Height
#     )
#     divider_line.line.color.rgb = MCKINSEY_COLORS["light_gray"]
#     divider_line.line.width = Pt(0.5)
def add_enhanced_text_and_image_slide(slide, bullets, image_path):
    """
    Creates an elegant side-by-side layout with bullet points on the left 
    and a large, professionally positioned image on the right.
    """

    # Left side - Text content
    text_width = Inches(5.5)
    text_height = Inches(4.5)

    tx_box = slide.shapes.add_textbox(
        Inches(0.8),
        Inches(1.8),
        text_width,
        text_height
    )
    tf = tx_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    # Loop through bullets
    for bullet in bullets:
        if isinstance(bullet, dict):
            # Main bullet point
            p = tf.add_paragraph()
            p.text = f"• {bullet.get('point','').strip()}"
            p.font.name = FONT_NAME
            p.font.size = Pt(16)
            p.font.color.rgb = MCKINSEY_COLORS["text"]
            p.level = 0
            p.space_after = Pt(6)

            # Sub description
            desc = bullet.get("desc", "").strip()
            if desc:
                desc_p = tf.add_paragraph()
                desc_p.text = desc
                desc_p.font.name = FONT_NAME
                desc_p.font.size = Pt(13)
                desc_p.font.color.rgb = MCKINSEY_COLORS["gray"]
                desc_p.level = 1
                desc_p.space_after = Pt(12)

        else:
            # Fallback if bullet is plain string
            p = tf.add_paragraph()
            p.text = f"• {str(bullet).strip()}"
            p.font.name = FONT_NAME
            p.font.size = Pt(16)
            p.font.color.rgb = MCKINSEY_COLORS["text"]
            p.level = 0
            p.space_after = Pt(12)

    # === Right side image handling (your existing code) ===
    img_stream = validated_image_bytes(image_path)
    if not img_stream:
        img_stream = validated_image_bytes("fallback.jpg")
    if img_stream:
        try:
            available_width = Inches(5.8)
            available_height = Inches(4.5)
            with Image.open(image_path if os.path.exists(image_path) else "fallback.jpg") as img:
                original_width, original_height = img.size
                aspect_ratio = original_width / original_height
                if aspect_ratio > (available_width.inches / available_height.inches):
                    img_width = available_width
                    img_height = Inches(available_width.inches / aspect_ratio)
                else:
                    img_height = available_height
                    img_width = Inches(available_height.inches * aspect_ratio)
                img_width = max(img_width, Inches(4.0))
                img_height = max(img_height, Inches(3.0))
            img_x = Inches(7.0)
            img_y = Inches(1.8) + (available_height - img_height) / 2
            pic = slide.shapes.add_picture(img_stream, img_x, img_y, width=img_width, height=img_height)
            pic.line.color.rgb = MCKINSEY_COLORS["light_gray"]
            pic.line.width = Pt(1)
        except Exception as e:
            print(f"[Image Error] {e}")


def calculate_dynamic_image_size(image_path, max_width=5.8, max_height=4.5):
    """
    Enhanced image size calculation for the new layout.
    """
    try:
        with Image.open(image_path) as img:
            original_width, original_height = img.size
            aspect_ratio = original_width / original_height
            
            # Calculate dimensions based on aspect ratio
            if aspect_ratio > (max_width / max_height):
                # Image is wider - constrain by width
                width = max_width
                height = max_width / aspect_ratio
            else:
                # Image is taller - constrain by height
                height = max_height
                width = max_height * aspect_ratio
            
            # Ensure minimum size for impact
            width = max(width, 4.0)
            height = max(height, 3.0)
            
            return Inches(width), Inches(height)
    except Exception as e:
        print(f"Error calculating image size: {e}")
        return Inches(5.0), Inches(3.5)  # Default fallback size
def get_slide_content_with_charts(topic, n_slides):
    prompt = f"""
    You are preparing a professional, data-driven PowerPoint presentation outline 
    with exactly {n_slides} slides on the topic "{topic}".

    Return ONLY valid JSON in this format:
    {{
      "intro": "<1-paragraph introduction>",
      "slides": [
        {{
          "title": "Slide Title",
          "insight": "One key insight",
          "type": "bullets",
          "data": ["point1", "point2", "point3"]
        }},
        {{
          "title": "Slide with Chart",
          "insight": "Key data-driven point",
          "type": "chart",
          "data": {{
            "type": "BAR",
            "data": [["Label1", 123], ["Label2", 456]],
            "source": "Source: Organization, 2025"
          }},
          "context": "Why this data matters"
        }}
      ]
    }}
    """
    return llm.invoke(prompt).content


def retry_and_fix_slides(parsed_slides, topic, n, max_retries=2):
    """
    FIXED: Replaces broken slides in their original positions instead of appending at the end.
    Ensures the total number of slides is 'n' while preserving slide sequence.
    """
    def is_broken(slide):
        data_str = str(slide.get("data", ""))
        return ("No content was provided" in data_str or 
                "No chart data was available" in data_str or
                not slide.get("data") or
                (isinstance(slide.get("data"), list) and len(slide.get("data")) == 0))

    retry_count = 0
    while retry_count < max_retries:
        # FIXED: Identify broken slides by their indices
        broken_indices = [i for i, slide in enumerate(parsed_slides["slides"]) if is_broken(slide)]
        missing_count = max(0, n - len(parsed_slides["slides"]))
        
        if not broken_indices and missing_count == 0:
            break  # All slides are good and we have the right count
        
        total_needed = len(broken_indices) + missing_count
        if total_needed == 0:
            break

        print(f"[Retry {retry_count + 1}] Fixing {len(broken_indices)} broken slides, generating {missing_count} additional slides...")

        # Generate new content for the needed slides
        new_raw = get_slide_content_with_charts(topic, total_needed)
        new_parsed = parse_mckinsey_response(new_raw)
        new_slides = new_parsed["slides"]

        # FIXED: Replace broken slides in their original positions
        replacement_index = 0
        for broken_idx in broken_indices:
            if replacement_index < len(new_slides):
                print(f"   Replacing broken slide at position {broken_idx + 1}")
                parsed_slides["slides"][broken_idx] = new_slides[replacement_index]
                replacement_index += 1

        # Add any remaining new slides for missing count
        for i in range(replacement_index, len(new_slides)):
            if len(parsed_slides["slides"]) < n:
                print(f"   Adding new slide at position {len(parsed_slides['slides']) + 1}")
                parsed_slides["slides"].append(new_slides[i])

        retry_count += 1

    # FIXED: Ensure we have exactly n slides, trim excess or pad with fallback
    current_count = len(parsed_slides["slides"])
    if current_count > n:
        print(f"   Trimming {current_count - n} excess slides")
        parsed_slides["slides"] = parsed_slides["slides"][:n]
    elif current_count < n:
        print(f"   Adding {n - current_count} fallback slides")
        for i in range(current_count, n):
            fallback_slide = {
                "title": f"Additional Insights {i + 1}",
                "insight": "Further analysis and strategic considerations",
                "type": "bullets",
                "data": [
                    "This area requires additional research and analysis",
                    "Strategic implications need further evaluation", 
                    "Recommend follow-up discussion with stakeholders"
                ],
                "context": ""
            }
            parsed_slides["slides"].append(fallback_slide)

    # Final validation: Replace any remaining broken slides with generic content
    for i, slide in enumerate(parsed_slides["slides"]):
        if is_broken(slide):
            print(f"   Final fallback for slide {i + 1}")
            parsed_slides["slides"][i] = {
                "title": f"Key Topic {i + 1}",
                "insight": "Strategic analysis and recommendations",
                "type": "bullets",
                "data": [
                    "Comprehensive analysis of current market conditions",
                    "Strategic recommendations for optimal outcomes",
                    "Implementation roadmap and next steps"
                ],
                "context": ""
            }

    return parsed_slides
def get_slide_content_with_charts(topic, n_slides):
    prompt = f"""
    You are preparing a professional, data-driven PowerPoint presentation outline 
    with exactly {n_slides} slides on the topic "{topic}".

    Rules:
    - Return ONLY valid JSON, no explanations or markdown.
    - "intro" must be maximum 1–2 lines (<=150 characters).
    - Each slide must include:
      - "title": max 8 words
      - "insight": one sentence
      - "type": either "bullets" or "chart"
      - If "type" = "bullets", "data" must be an array of objects:
        {{"point": "short phrase", "desc": "1–2 line explanation"}}
      - If "type" = "chart", "data" must be JSON:
        {{
          "type": "BAR" | "LINE" | "PIE" | "COLUMN" | "DOUGHNUT" | "AREA" | "SCATTER" | "STACKED_BAR",
          "data": [["Label1", 123], ["Label2", 456]],
          "source": "Source: Organization, 2025"
        }}
        and optionally include "context": "Why this data matters"

    JSON format:
    {{
      "intro": "short intro text",
      "slides": [
        {{
          "title": "Slide Title",
          "insight": "Key insight",
          "type": "bullets",
          "data": [
            {{"point": "Clarity", "desc": "Use straightforward language to avoid confusion"}},
            {{"point": "Audience Awareness", "desc": "Adapt tone and style to the readers"}}
          ]
        }},
        {{
          "title": "Slide with Chart",
          "insight": "Data-driven point",
          "type": "chart",
          "data": {{
            "type": "COLUMN",
            "data": [["X", 10], ["Y", 20]],
            "source": "Source: Example Org, 2025"
          }},
          "context": "Why this data is important"
        }}
      ]
    }}
    """
    return llm.invoke(prompt).content

if __name__ == "__main__":
    mode = input("Choose input mode (1=Topic, 2=Paragraph): ").strip()

    
    def parse_json_slides(text: str):
        try:
            # 1. Strip markdown fences if present
            cleaned = text.strip()
            if cleaned.startswith("```"):
                cleaned = re.sub(r"^```(json)?", "", cleaned, flags=re.IGNORECASE).strip()
                cleaned = re.sub(r"```$", "", cleaned).strip()

            # 2. Extract JSON object if there's extra text around it
            match = re.search(r"\{.*\}", cleaned, re.DOTALL)
            if match:
                cleaned = match.group(0)

            # 3. Parse JSON
            return json.loads(cleaned)

        except json.JSONDecodeError as e:
            print(f"[Error] JSON parse failed: {e}")
            print(f"[Debug Raw Start] {text[:300]}...")
            return {"intro": "", "slides": []}

    if mode == "1":
        # === Topic-based workflow (JSON) ===
        topic = input("Enter your presentation topic: ").strip()
        try:
            n = int(input("Number of slides: "))
            if n < 1:
                raise ValueError("Slide count must be at least 1.")
        except Exception as e:
            print(f"Invalid input: {e}")
            exit()

        print("\n[1/3] Generating JSON content with AI...")
        raw_content = get_slide_content_with_charts(topic, n)

        print("[2/3] Parsing JSON response...")
        parsed_slides = parse_json_slides(raw_content)

        print(f"[3/3] Building '{topic}' presentation with {len(parsed_slides['slides'])} slides...")
        if not parsed_slides["slides"]:
            print("\nFailed to generate or parse JSON content. Please try again.")
        else:
            build_mckinsey_ppt(parsed_slides, topic)

    elif mode == "2":
        # === Paragraph-based workflow (JSON) ===
        context_text = input("Enter paragraph/context for your presentation: ").strip()

        print("\n[Step 1] Detecting content type...")
        category = classify_paragraph_type(context_text)
        print(f"→ Detected category: {category}")

        print("\n[Step 2] Refining input text...")
        refined_text = refine_paragraph_input(context_text, category)
        print("\n[Refined Version Used for Slides]:\n")
        print(refined_text, "\n")

        def generate_topic_from_paragraph(context_text):
            prompt = f"""
            Generate a concise, professional presentation title (max 6 words) from this content:
            "{context_text[:200]}..."
            
            Return only the title, no quotes or formatting.
            """
            return llm.invoke(prompt).content.strip()

        topic = generate_topic_from_paragraph(context_text)

        print("\n[Step 3] Generating JSON slides...")
        raw_content = get_slide_content_from_paragraph(refined_text, category)

        print("[Step 4] Parsing JSON response...")
        parsed_slides = parse_json_slides(raw_content)

        print(f"[Step 5] Building '{topic}' presentation with {len(parsed_slides['slides'])} slides...")
        if not parsed_slides["slides"]:
            print("\nFailed to generate or parse JSON content. Please try again.")
        else:
            build_mckinsey_ppt(parsed_slides, topic)

    else:
        print("Invalid mode. Exiting.")
        exit()
