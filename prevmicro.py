import os
import re
import subprocess
import wave

import tempfile
import shutil
import logging
from pathlib import Path
import requests
import json
from pptx.enum.text import MSO_AUTO_SIZE
import re
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_DATA_LABEL_POSITION
from pptx.chart.data import CategoryChartData
from dotenv import load_dotenv
from langchain_groq import ChatGroq
from PIL import Image, ImageDraw
from io import BytesIO
from bs4 import BeautifulSoup
from datetime import datetime
import cv2
import numpy as np
from gtts import gTTS
import moviepy as mp

import tempfile
import shutil
from pathlib import Path
import time
import threading
from concurrent.futures import ThreadPoolExecutor
import logging

# === Setup logging ===
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# === Load environment variables ===
load_dotenv()
GROQ_API_KEY = os.getenv("GROQ_API_KEY")
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY_1")
GOOGLE_CSE_ID = os.getenv("GOOGLE_CSE_ID_1")
UNSPLASH_ACCESS_KEY = os.getenv("UNSPLASH_ACCESS_KEY")

# === Initialize LLM ===
llm = ChatGroq(
    groq_api_key=GROQ_API_KEY,
    model_name="openai/gpt-oss-20b",
)

# === Enhanced McKinsey Style Constants ===
MCKINSEY_COLORS = {
    "blue": RGBColor(12, 74, 126),
    "light_blue": RGBColor(79, 129, 189),
    "dark_blue": RGBColor(6, 45, 85),
    "accent_blue": RGBColor(102, 170, 238),
    "gray": RGBColor(89, 89, 89),
    "light_gray": RGBColor(217, 217, 217),
    "dark_gray": RGBColor(64, 64, 64),
    "background": RGBColor(255, 255, 255),
    "text": RGBColor(0, 0, 0),
    "white": RGBColor(255, 255, 255)
}
FONT_NAME = "Aptos Display"

# === Video Configuration ===
VIDEO_CONFIG = {
    "width": 1920,
    "height": 1080,
    "fps": 30,
    "slide_duration": 8,  # seconds per slide (will be adjusted based on audio)
    "transition_duration": 0.5,  # seconds for transitions
    "background_color": (255, 255, 255),  # White background
    "output_format": "mp4"
}

def get_slide_content_with_charts(topic, n_slides):
    prompt = f"""
    You are preparing a professional, data-driven PowerPoint presentation outline 
    with exactly {n_slides} slides on the topic "{topic}".

    Rules:
    - Return ONLY valid JSON, no explanations or markdown.
    - "intro" must be maximum 1â€“2 lines (<=150 characters).
    - Each slide must include:
      - "title": max 8 words
      - "insight": one sentence
      - "type": either "bullets" or "chart"
      - If "type" = "bullets", "data" must be an array of objects:
        {{"point": "short phrase", "desc": "1â€“2 line explanation"}}
      - If "type" = "chart", "data" must be JSON:
        {{
          "type": "BAR" | "LINE" | "PIE" | "COLUMN" | "DOUGHNUT" | "AREA" | "SCATTER" | "STACKED_BAR",
          "data": [["Label1", 123], ["Label2", 456]],
          "source": "Source: Organization, 2025"
        }}
        and optionally include "context": "Why this data matters"

    JSON format:
    {{
      "intro": "short intro text",
      "slides": [
        {{
          "title": "Slide Title",
          "insight": "Key insight",
          "type": "bullets",
          "data": [
            {{"point": "Clarity", "desc": "Use straightforward language to avoid confusion"}},
            {{"point": "Audience Awareness", "desc": "Adapt tone and style to the readers"}}
          ]
        }},
        {{
          "title": "Slide with Chart",
          "insight": "Data-driven point",
          "type": "chart",
          "data": {{
            "type": "COLUMN",
            "data": [["X", 10], ["Y", 20]],
            "source": "Source: Example Org, 2025"
          }},
          "context": "Why this data is important"
        }}
      ]
    }}
    """
    return llm.invoke(prompt).content

def parse_mckinsey_response(text):
    slides = []
    intro_match = re.search(r"Introduction:\s*(.*?)\n---", text, re.DOTALL)
    introduction = ""

    slide_chunks = text.split('---')
    for chunk in slide_chunks:
        if not chunk.strip():
            continue

        try:
            title_match = re.search(r"Title: (.*?)\n", chunk, re.IGNORECASE)
            insight_match = re.search(r"Key Insight: (.*?)\n", chunk, re.IGNORECASE)
            context_match = re.search(r"Context: (.*?)\n", chunk, re.IGNORECASE)

            title = title_match.group(1).strip() if title_match else f"Slide {len(slides)+1}"
            insight = insight_match.group(1).strip() if insight_match else None
            context = context_match.group(1).strip() if context_match else ""

            if not insight:
                print(f"Skipping slide due to missing insight: {chunk[:100]}")
                continue

            content = {"title": title, "insight": insight, "context": context}

            chart_match = re.search(r"Chart:\s*```json\n(.*?)\n```", chunk, re.DOTALL)
            if chart_match:
                try:
                    chart_data = json.loads(chart_match.group(1).strip())
                    content["type"] = "chart"
                    content["data"] = chart_data
                except json.JSONDecodeError:
                    print(f"[Warning] Chart JSON malformed. Falling back to bullets.")
                    content["type"] = "bullets"
                    content["data"] = ["No chart data was available."]
            else:
                bullets_match = re.search(r"Bullets:\s*(.*)", chunk, re.DOTALL)
                bullets_raw = bullets_match.group(1) if bullets_match else ""
                bullets = [b.strip("- ").strip() for b in bullets_raw.split('\n') if b.strip()]
                if not bullets:
                    bullets = ["No content was provided by the AI for this slide."]
                content["type"] = "bullets"
                content["data"] = bullets

            slides.append(content)
        except Exception as e:
            print(f"Could not parse a slide chunk, skipping. Error: {e}\nChunk: {chunk[:100]}...")
            continue

    return {"intro": introduction, "slides": slides}

def parse_json_slides(text: str):
    try:
        # 1. Strip markdown fences if present
        cleaned = text.strip()
        if cleaned.startswith("```"):
            cleaned = re.sub(r"^```(json)?", "", cleaned, flags=re.IGNORECASE).strip()
            cleaned = re.sub(r"```$", "", cleaned).strip()

        # 2. Extract JSON object if there's extra text around it
        match = re.search(r"\{.*\}", cleaned, re.DOTALL)
        if match:
            cleaned = match.group(0)

        # 3. Parse JSON
        return json.loads(cleaned)

    except json.JSONDecodeError as e:
        print(f"[Error] JSON parse failed: {e}")
        print(f"[Debug Raw Start] {text[:300]}...")
        return {"intro": "", "slides": []}

def add_enhanced_title_slide(prs, topic, intro):
    """Creates a stunning McKinsey-style title slide with professional layout and design elements."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = MCKINSEY_COLORS["background"]

    # Background geometric elements for visual appeal
    bg_rect = slide.shapes.add_shape(1, Inches(8), Inches(0), Inches(3.33), Inches(7.5))
    bg_rect.fill.solid()
    bg_rect.fill.fore_color.rgb = MCKINSEY_COLORS["light_blue"]
    bg_rect.line.fill.background()
    
    # Diagonal accent shape
    accent_shape = slide.shapes.add_shape(1, Inches(0), Inches(6.7), Inches(13.33), Inches(1.0))
    accent_shape.fill.solid()
    accent_shape.fill.fore_color.rgb = MCKINSEY_COLORS["blue"]
    accent_shape.line.fill.background()

    # Main Title - Bold and prominent
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(10), Inches(2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = topic.upper()
    p.font.name = FONT_NAME
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = MCKINSEY_COLORS["dark_blue"]
    p.alignment = PP_ALIGN.LEFT

    # Professional subtitle with current date
    current_date = datetime.now().strftime("%B %Y")
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(10), Inches(0.8))
    tf2 = subtitle_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = f"Strategic Analysis & Business Intelligence | {current_date}"
    p2.font.name = FONT_NAME
    p2.font.size = Pt(20)
    p2.font.color.rgb = MCKINSEY_COLORS["gray"]
    p2.alignment = PP_ALIGN.LEFT

    # Executive Summary Box with better formatting
    intro_box = slide.shapes.add_textbox(Inches(1), Inches(4.0), Inches(10), Inches(1.5))
    tf3 = intro_box.text_frame
    tf3.word_wrap = True
    tf3.auto_size = MSO_AUTO_SIZE.NONE
    tf3.margin_left = Inches(0.2)
    tf3.margin_right = Inches(0.2)
    tf3.margin_top = Inches(0.1)
    tf3.margin_bottom = Inches(0.1)
    
    p3 = tf3.paragraphs[0]
    p3.font.name = FONT_NAME
    p3.font.size = Pt(14)
    p3.font.bold = True
    p3.font.color.rgb = MCKINSEY_COLORS["blue"]
    p3.alignment = PP_ALIGN.LEFT
    
    # Add the actual introduction content
    p4 = tf3.add_paragraph()
    p4.text = intro
    p4.font.name = FONT_NAME
    p4.font.size = Pt(16)
    p4.font.color.rgb = MCKINSEY_COLORS["text"]
    p4.alignment = PP_ALIGN.LEFT
    p4.space_before = Pt(8)

    # Professional branding footer
    brand_box = slide.shapes.add_textbox(Inches(1), Inches(7), Inches(6), Inches(0.4))
    tf4 = brand_box.text_frame
    p5 = tf4.paragraphs[0]
    p5.text = "MADE BY ENTHRAL AI"
    p5.font.name = FONT_NAME
    p5.font.size = Pt(10)
    p5.font.italic = True
    p5.font.color.rgb = MCKINSEY_COLORS["white"]
    p5.alignment = PP_ALIGN.LEFT

    # Decorative elements
    line1 = slide.shapes.add_shape(1, Inches(1), Inches(4.2), Inches(4), Inches(0))
    line1.line.color.rgb = MCKINSEY_COLORS["accent_blue"]
    line1.line.width = Pt(4)
    
    line2 = slide.shapes.add_shape(1, Inches(1), Inches(4.3), Inches(2), Inches(0))
    line2.line.color.rgb = MCKINSEY_COLORS["blue"]
    line2.line.width = Pt(2)

def add_chart_slide_with_context(slide, chart_info, context):
    if context:
        tx_box = slide.shapes.add_textbox(Inches(1), Inches(1.6), Inches(11), Inches(0.6))
        tf = tx_box.text_frame
        p = tf.paragraphs[0]
        p.text = context
        p.font.size = Pt(14)
        p.font.name = FONT_NAME
        p.font.color.rgb = MCKINSEY_COLORS["gray"]

    add_chart_slide(slide, chart_info)

def calculate_dynamic_image_size(image_path, max_width=5.8, max_height=4.5):
    """Enhanced image size calculation for the new layout."""
    try:
        with Image.open(image_path) as img:
            original_width, original_height = img.size
            aspect_ratio = original_width / original_height
            
            # Calculate dimensions based on aspect ratio
            if aspect_ratio > (max_width / max_height):
                width = max_width
                height = max_width / aspect_ratio
            else:
                height = max_height
                width = max_height * aspect_ratio
            
            # Ensure minimum size for impact
            width = max(width, 4.0)
            height = max(height, 3.0)
            
            return Inches(width), Inches(height)
    except Exception as e:
        print(f"Error calculating image size: {e}")
        return Inches(5.0), Inches(3.5)

def build_mckinsey_ppt(parsed_data, topic):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    add_enhanced_title_slide(prs, topic, parsed_data["intro"])

    for i, slide_content in enumerate(parsed_data["slides"], start=1):
        slide = prs.slides.add_slide(blank_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = MCKINSEY_COLORS["background"]

        add_enhanced_header(slide, slide_content['title'], slide_content['insight'])
        add_enhanced_footer(slide, i + 1)

        if slide_content['type'] == 'chart':
            add_chart_slide_with_context(slide, slide_content['data'], slide_content.get("context", ""))
        else:
            image_path = fetch_image(f"{slide_content['title']} {slide_content['insight']}", topic)
            add_enhanced_text_and_image_slide(slide, slide_content['data'], image_path)

    filename = topic.strip().replace(" ", "_") + "_McKinsey_Style.pptx"
    prs.save(filename)
    print(f"\nPresentation saved as: {filename}")
    return filename

def add_enhanced_header(slide, title, insight):
    """Adds a professionally styled title and key insight with better visual hierarchy."""
    # Background header area
    header_bg = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.4))
    header_bg.fill.solid()
    header_bg.fill.fore_color.rgb = MCKINSEY_COLORS["background"]
    header_bg.line.color.rgb = MCKINSEY_COLORS["light_gray"]
    header_bg.line.width = Pt(1)
    
    # Title
    tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.65))
    tf = tx_box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = title.upper()
    p.font.name = FONT_NAME
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = MCKINSEY_COLORS["dark_blue"]

    # Key Insight with better styling
    insight_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(12), Inches(0.5))
    tf2 = insight_box.text_frame
    tf2.margin_left = Inches(0.1)
    p_insight = tf2.paragraphs[0]
    p_insight.text = f"Key Insight: {insight}"
    p_insight.font.name = FONT_NAME
    p_insight.font.size = Pt(14)
    p_insight.font.color.rgb = MCKINSEY_COLORS["gray"]
    p_insight.font.italic = True

    # Enhanced blue accent line
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.35), Inches(12), Inches(0))
    line.line.color.rgb = MCKINSEY_COLORS["blue"]
    line.line.width = Pt(3)

def add_enhanced_footer(slide, slide_number):
    """Adds a professionally styled footer with slide number and branding."""
    # Footer background
    footer_bg = slide.shapes.add_shape(1, Inches(0), Inches(6.9), Inches(13.33), Inches(0.6))
    footer_bg.fill.solid()
    footer_bg.fill.fore_color.rgb = MCKINSEY_COLORS["light_gray"]
    footer_bg.line.fill.background()
    
    # Slide number
    tx_box = slide.shapes.add_textbox(Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.4))
    tf = tx_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{slide_number}"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FONT_NAME
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = MCKINSEY_COLORS["blue"]
    
    # Professional footer text
    footer_text_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(6), Inches(0.4))
    tf2 = footer_text_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "MADE BY ENTHRAL AI"
    p2.font.name = FONT_NAME
    p2.font.size = Pt(9)
    p2.font.color.rgb = MCKINSEY_COLORS["dark_gray"]

def fetch_image(prompt: str, topic: str) -> str:
    """Enhanced image fetching with multiple sources."""
    from urllib.parse import quote
    
    headers = {"User-Agent": "Mozilla/5.0"}
    query = f"{topic} {prompt}".strip()
    encoded = quote(query)

    def is_vertical(img):
        w, h = img.size
        aspect_ratio = w / h
        return aspect_ratio < 1.25

    def try_save_image(url, label, require_vertical=True):
        try:
            r = requests.get(url, stream=True, headers=headers, timeout=10)
            if r.status_code == 200 and 'image' in r.headers.get('Content-Type', ''):
                img = Image.open(BytesIO(r.content)).convert("RGB")
                if not require_vertical or is_vertical(img):
                    filename = f"{label}_{re.sub(r'[^a-zA-Z0-9]', '_', url)[:30]}.jpg"
                    img.save(filename, format="JPEG")
                    return filename
        except:
            pass
        return None

    def fetch_from_google(require_vertical=True):
        try:
            if not GOOGLE_API_KEY or not GOOGLE_CSE_ID:
                return None
                
            g_url = f"https://www.googleapis.com/customsearch/v1?q={encoded}&searchType=image&num=8&key={GOOGLE_API_KEY}&cx={GOOGLE_CSE_ID}"
            response = requests.get(g_url, timeout=15)
            
            if response.status_code != 200:
                return None
                
            data = response.json()
            
            if 'items' not in data:
                return None
                
            for i, item in enumerate(data['items']):
                url = item.get("link", "")
                if url:
                    img_path = try_save_image(url, f"google_{i}", require_vertical)
                    if img_path:
                        return img_path
                        
        except Exception as e:
            pass
        return None

    def create_fallback():
        try:
            img = Image.new('RGB', (800, 600), color=(245, 245, 245))
            draw = ImageDraw.Draw(img)
            draw.rectangle([40, 40, 760, 560], outline=(12, 74, 126), width=4)
            draw.rectangle([80, 80, 720, 520], outline=(217, 217, 217), width=2)
            
            try:
                from PIL import ImageFont
                font_large = ImageFont.load_default()
                font_small = ImageFont.load_default()
            except:
                font_large = ImageFont.load_default()
                font_small = ImageFont.load_default()
            
            draw.text((400, 250), "Professional", fill=(12, 74, 126), anchor='mm', font=font_large)
            draw.text((400, 300), "Image Placeholder", fill=(12, 74, 126), anchor='mm', font=font_large)
            draw.text((400, 380), f"Topic: {topic}", fill=(89, 89, 89), anchor='mm', font=font_small)
            draw.ellipse([320, 150, 480, 200], outline=(79, 129, 189), width=3)
            
            img.save('fallback.jpg', 'JPEG', quality=90)
            return True
        except:
            return False

    # Try different methods
    for orientation_pref in [True, False]:
        image_path = fetch_from_google(require_vertical=orientation_pref)
        if image_path:
            return image_path
    
    # Fallback
    if not os.path.exists("fallback.jpg"):
        create_fallback()
    
    return "fallback.jpg" if os.path.exists("fallback.jpg") else None

def validated_image_bytes(path, max_width=800, max_height=500):
    """Enhanced image validation with better error handling and optimization."""
    try:
        if not path or not os.path.exists(path):
            return None
            
        with Image.open(path) as img:
            img.verify()
            
        with Image.open(path) as img:
            if img.mode != 'RGB':
                img = img.convert("RGB")
            
            original_w, original_h = img.size
            img.thumbnail((max_width, max_height), Image.Resampling.LANCZOS)
            
            stream = BytesIO()
            img.save(stream, format="PNG", optimize=True, quality=90)
            stream.seek(0)
            
            return stream
            
    except Exception as e:
        return None

def add_chart_slide(slide, chart_info):
    """Adds a dynamically generated chart to the slide."""
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_DATA_LABEL_POSITION

    try:
        raw_data = chart_info.get("data", [])
        if not raw_data or not isinstance(raw_data, list):
            raise ValueError("Chart data is missing or malformed.")

        chart_type_str = chart_info.get("type", "BAR").upper()
        chart_type_map = {
            "BAR": XL_CHART_TYPE.BAR_CLUSTERED,
            "COLUMN": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "LINE": XL_CHART_TYPE.LINE,
            "PIE": XL_CHART_TYPE.PIE,
            "DOUGHNUT": XL_CHART_TYPE.DOUGHNUT,
            "AREA": XL_CHART_TYPE.AREA,
            "SCATTER": XL_CHART_TYPE.XY_SCATTER_LINES,
            "STACKED_BAR": XL_CHART_TYPE.BAR_STACKED
        }
        chart_type = chart_type_map.get(chart_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED)

        chart_data = CategoryChartData()

        # Handle different data formats
        if all(isinstance(row, list) and len(row) > 2 for row in raw_data):
            categories = [row[0] for row in raw_data]
            num_series = len(raw_data[0]) - 1
            series_data = [[] for _ in range(num_series)]

            for row in raw_data:
                for i in range(num_series):
                    val = row[i + 1]
                    series_data[i].append(float(val) if isinstance(val, (int, float)) else 0)

            chart_data.categories = categories
            for i, values in enumerate(series_data):
                chart_data.add_series(f"Series {i + 1}", values)

        elif all(isinstance(row, list) and len(row) == 2 for row in raw_data):
            chart_data.categories = [str(row[0]) for row in raw_data]
            values = [float(row[1]) if isinstance(row[1], (int, float)) else 0.0 for row in raw_data]
            chart_data.add_series("", values)

        else:
            raise ValueError("Invalid chart data format.")

        # Create chart
        x, y, cx, cy = Inches(1), Inches(1.7), Inches(11), Inches(4.5)
        graphic_frame = slide.shapes.add_chart(chart_type, x, y, cx, cy, chart_data)
        chart = graphic_frame.chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False

        plot = chart.plots[0]
        plot.has_data_labels = True
        plot.data_labels.font.size = Pt(12)
        plot.data_labels.font.bold = True
        plot.data_labels.font.color.rgb = MCKINSEY_COLORS["gray"]

        # Chart source
        source_text = chart_info.get("source", "")
        if source_text:
            tx_box = slide.shapes.add_textbox(Inches(1), Inches(6.4), Inches(11), Inches(0.4))
            p = tx_box.text_frame.paragraphs[0]
            p.text = source_text
            p.font.size = Pt(9)
            p.font.italic = True
            p.font.color.rgb = MCKINSEY_COLORS["gray"]

    except Exception as e:
        print(f"Failed to create chart: {e}")
        tx_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(1))
        tx_box.text_frame.paragraphs[0].text = "Error: Could not generate the requested chart."

def add_enhanced_text_and_image_slide(slide, bullets, image_path):
    """Creates an elegant side-by-side layout with bullet points and image."""
    text_width = Inches(5.5)
    text_height = Inches(4.5)

    tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), text_width, text_height)
    tf = tx_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    for bullet in bullets:
        if isinstance(bullet, dict):
            p = tf.add_paragraph()
            p.text = f"â€¢ {bullet.get('point','').strip()}"
            p.font.name = FONT_NAME
            p.font.size = Pt(16)
            p.font.color.rgb = MCKINSEY_COLORS["text"]
            p.level = 0
            p.space_after = Pt(6)

            desc = bullet.get("desc", "").strip()
            if desc:
                desc_p = tf.add_paragraph()
                desc_p.text = desc
                desc_p.font.name = FONT_NAME
                desc_p.font.size = Pt(13)
                desc_p.font.color.rgb = MCKINSEY_COLORS["gray"]
                desc_p.level = 1
                desc_p.space_after = Pt(12)
        else:
            p = tf.add_paragraph()
            p.text = f"â€¢ {str(bullet).strip()}"
            p.font.name = FONT_NAME
            p.font.size = Pt(16)
            p.font.color.rgb = MCKINSEY_COLORS["text"]
            p.level = 0
            p.space_after = Pt(12)

    # Add image
    img_stream = validated_image_bytes(image_path)
    if not img_stream:
        img_stream = validated_image_bytes("fallback.jpg")

    if img_stream:
        try:
            available_width = Inches(5.8)
            available_height = Inches(4.5)
            with Image.open(image_path if os.path.exists(image_path) else "fallback.jpg") as img:
                original_width, original_height = img.size
                aspect_ratio = original_width / original_height
                if aspect_ratio > (available_width.inches / available_height.inches):
                    img_width = available_width
                    img_height = Inches(available_width.inches / aspect_ratio)
                else:
                    img_height = available_height
                    img_width = Inches(available_height.inches * aspect_ratio)
                img_width = max(img_width, Inches(4.0))
                img_height = max(img_height, Inches(3.0))
            img_x = Inches(7.0)
            img_y = Inches(1.8) + (available_height - img_height) / 2
            pic = slide.shapes.add_picture(img_stream, img_x, img_y, width=img_width, height=img_height)
            pic.line.color.rgb = MCKINSEY_COLORS["light_gray"]
            pic.line.width = Pt(1)
        except Exception as e:
            print(f"[Image Error] {e}")

# === AI NARRATION GENERATION ===

# def generate_narration_script(parsed_data, topic):
#     """Generate AI narration script for the entire presentation."""
#     logger.info("Generating narration script with AI...")
    
#     prompt = f"""
#     Create a professional narration script for a PowerPoint presentation about "{topic}".
#     The script should be engaging, clear, and suitable for text-to-speech conversion.
    
#     Presentation data:
#     Introduction: {parsed_data.get("intro", "")}
    
#     Slides:
#     """
    
#     for i, slide in enumerate(parsed_data.get("slides", []), 1):
#         prompt += f"\nSlide {i}: {slide.get('title', '')}\n"
#         prompt += f"Key Insight: {slide.get('insight', '')}\n"
        
#         if slide.get('type') == 'chart':
#             chart_data = slide.get('data', {})
#             prompt += f"Chart Type: {chart_data.get('type', 'Unknown')}\n"
#             prompt += f"Chart Source: {chart_data.get('source', 'No source')}\n"
#             if 'context' in slide:
#                 prompt += f"Context: {slide['context']}\n"
#         else:
#             bullets = slide.get('data', [])
#             if isinstance(bullets, list):
#                 for bullet in bullets:
#                     if isinstance(bullet, dict):
#                         prompt += f"- {bullet.get('point', '')}: {bullet.get('desc', '')}\n"
#                     else:
#                         prompt += f"- {str(bullet)}\n"
    
#     prompt += f"""
    
#     Create a narration script with the following structure:
#     1. Title slide: Welcome and introduction (20-30 seconds)
#     2. For each content slide: Explain the title, key insight, and main points (30-45 seconds each)
#     3. Keep the language conversational but professional
#     4. Use smooth transitions between slides
#     5. End with a brief conclusion
    
#     Format the response as JSON:
#     {{
#         "title_narration": "Welcome text for title slide...",
#         "slide_narrations": [
#             "Narration for slide 1...",
#             "Narration for slide 2...",
#             ...
#         ],
#         "conclusion": "Brief closing remarks..."
#     }}
#     """
    
#     try:
#         response = llm.invoke(prompt).content
#         # Parse JSON response
#         narration_data = json.loads(response.strip())
#         return narration_data
#     except Exception as e:
#         logger.error(f"Failed to generate narration script: {e}")
#         # Fallback narration
#         return {
#             "title_narration": f"Welcome to this presentation about {topic}. Let's explore the key insights and analysis.",
#             "slide_narrations": [f"This slide covers {slide.get('title', f'slide {i+1}')}. {slide.get('insight', 'Key information is presented here.')}" 
#                                for i, slide in enumerate(parsed_data.get("slides", []))],
#             "conclusion": "Thank you for your attention. This concludes our presentation."
#         }
def generate_narration_script(parsed_data, topic):
    """Generate AI narration script for the entire presentation with better error handling."""
    logger.info("Generating narration script with AI...")
    
    try:
        prompt = f"""
        Create a professional narration script for a PowerPoint presentation about "{topic}".
        The script should be engaging, clear, and suitable for text-to-speech conversion.
        
        Presentation data:
        Introduction: {parsed_data.get("intro", "")}
        
        Slides:
        """
        
        for i, slide in enumerate(parsed_data.get("slides", []), 1):
            prompt += f"\nSlide {i}: {slide.get('title', '')}\n"
            prompt += f"Key Insight: {slide.get('insight', '')}\n"
            
            if slide.get('type') == 'chart':
                chart_data = slide.get('data', {})
                prompt += f"Chart Type: {chart_data.get('type', 'Unknown')}\n"
                prompt += f"Chart Source: {chart_data.get('source', 'No source')}\n"
                if 'context' in slide:
                    prompt += f"Context: {slide['context']}\n"
            else:
                bullets = slide.get('data', [])
                if isinstance(bullets, list):
                    for bullet in bullets[:3]:  # Limit to first 3 bullets for prompt
                        if isinstance(bullet, dict):
                            prompt += f"- {bullet.get('point', '')}: {bullet.get('desc', '')}\n"
                        else:
                            prompt += f"- {str(bullet)}\n"
        
        prompt += f"""
        
        Create a narration script with the following structure:
        1. Title slide: Welcome and introduction (20-30 seconds)
        2. For each content slide: Explain the title, key insight, and main points (30-45 seconds each)
        3. Keep the language conversational but professional
        4. Use smooth transitions between slides
        5. End with a brief conclusion
        
        Return only a JSON object in this exact format:
        {{
            "title_narration": "Welcome text for title slide...",
            "slide_narrations": [
                "Narration for slide 1...",
                "Narration for slide 2..."
            ],
            "conclusion": "Brief closing remarks..."
        }}
        """
        
        response = llm.invoke(prompt).content.strip()
        
        # Clean up response - remove markdown formatting if present
        if response.startswith("```json"):
            response = response.replace("```json", "").replace("```", "").strip()
        elif response.startswith("```"):
            response = response.replace("```", "", 1).replace("```", "").strip()
        
        # Parse JSON response
        import json
        try:
            narration_data = json.loads(response)
            
            # Validate the structure
            if not isinstance(narration_data, dict):
                raise ValueError("Response is not a dictionary")
            
            if "title_narration" not in narration_data:
                narration_data["title_narration"] = f"Welcome to this presentation about {topic}."
            
            if "slide_narrations" not in narration_data or not isinstance(narration_data["slide_narrations"], list):
                narration_data["slide_narrations"] = []
            
            if "conclusion" not in narration_data:
                narration_data["conclusion"] = "Thank you for your attention."
            
            return narration_data
            
        except json.JSONDecodeError as e:
            logger.error(f"JSON decode error: {e}")
            logger.error(f"Response was: {response[:500]}...")
            raise
            
    except Exception as e:
        logger.error(f"Failed to generate narration script: {e}")
        # Fallback narration
        return {
            "title_narration": f"Welcome to this presentation about {topic}. Let's explore the key insights and analysis.",
            "slide_narrations": [f"This slide covers {slide.get('title', f'slide {i+1}')}. {slide.get('insight', 'Key information is presented here.')}" 
                               for i, slide in enumerate(parsed_data.get("slides", []))],
            "conclusion": "Thank you for your attention. This concludes our presentation."
        }
def create_audio_from_text(text, output_path, lang='en', slow=False):
    """Create audio file from text using gTTS."""
    try:
        logger.info(f"Creating audio: {output_path[:50]}...")
        tts = gTTS(text=text, lang=lang, slow=slow)
        tts.save(output_path)
        return output_path
    except Exception as e:
        logger.error(f"Failed to create audio: {e}")
        return None

# def get_audio_duration(audio_path):
#     """Get the duration of an audio file."""
#     try:
#         audio = mp.AudioFileClip(audio_path)
#         duration = audio.duration
#         audio.close()
#         return duration
#     except Exception as e:
#         logger.error(f"Failed to get audio duration: {e}")
#         return VIDEO_CONFIG["slide_duration"]  # fallback duration

def convert_mp3_to_wav(mp3_path):
    """Convert MP3 to WAV using ffmpeg if available."""
    wav_path = mp3_path.replace('.mp3', '.wav')
    try:
        # Try using ffmpeg if available
        if shutil.which('ffmpeg'):
            subprocess.run([
                'ffmpeg', '-i', mp3_path, '-acodec', 'pcm_s16le', 
                '-ar', '44100', wav_path, '-y'
            ], check=True, capture_output=True, text=True)
            return wav_path
        else:
            logger.warning("ffmpeg not found, using original MP3 file")
            return mp3_path
    except subprocess.CalledProcessError as e:
        logger.warning(f"Failed to convert MP3 to WAV: {e}, using original file")
        return mp3_path
    except Exception as e:
        logger.error(f"Audio conversion error: {e}")
        return mp3_path
# def create_video_from_slides_and_audio(image_files, audio_files, output_path):
#     """Create video from slide images and audio files using OpenCV."""
#     logger.info("Creating video from slides and audio using OpenCV...")
    
#     try:
#         # Video writer setup
#         fourcc = cv2.VideoWriter_fourcc(*'mp4v')
#         video_writer = cv2.VideoWriter(
#             output_path, 
#             fourcc, 
#             VIDEO_CONFIG["fps"], 
#             (VIDEO_CONFIG["width"], VIDEO_CONFIG["height"])
#         )
        
#         temp_audio_files = []
#         total_duration = 0
        
#         for i, (image_path, audio_path) in enumerate(zip(image_files, audio_files)):
#             if not os.path.exists(image_path) or not os.path.exists(audio_path):
#                 logger.warning(f"Missing file for slide {i+1}")
#                 continue
            
#             # Get audio duration
#             audio_duration = get_audio_duration(audio_path)
#             total_duration += audio_duration
#             logger.info(f"Slide {i+1}: {audio_duration:.1f}s duration")
            
#             # Convert MP3 to WAV if needed for better compatibility
#             wav_path = convert_mp3_to_wav(audio_path)
#             temp_audio_files.append((wav_path, audio_path))  # Store both paths
            
#             # Load and process image
#             img = cv2.imread(image_path)
#             if img is None:
#                 logger.error(f"Could not load image: {image_path}")
#                 continue
                
#             # Resize image to video dimensions
#             img_resized = cv2.resize(img, (VIDEO_CONFIG["width"], VIDEO_CONFIG["height"]))
            
#             # Calculate number of frames for this slide
#             num_frames = int(audio_duration * VIDEO_CONFIG["fps"])
            
#             # Write frames for this slide
#             for frame_num in range(num_frames):
#                 # Add simple fade effect for transitions (optional)
#                 if i > 0 and frame_num < int(VIDEO_CONFIG["fps"] * VIDEO_CONFIG["transition_duration"]):
#                     # Fade in effect
#                     alpha = frame_num / (VIDEO_CONFIG["fps"] * VIDEO_CONFIG["transition_duration"])
#                     # Create a simple fade by blending with a black frame
#                     black_frame = np.zeros_like(img_resized)
#                     faded_img = cv2.addWeighted(img_resized, alpha, black_frame, 1-alpha, 0)
#                     video_writer.write(faded_img)
#                 else:
#                     video_writer.write(img_resized)
            
#             logger.info(f"Processed slide {i+1}: {num_frames} frames ({audio_duration:.1f}s)")
        
#         video_writer.release()
#         logger.info(f"Video frames written. Total duration: {total_duration:.1f}s")
        
#         # Now we need to add audio to the video
#         final_output_path = output_path.replace('.mp4', '_final.mp4')
        
#         if shutil.which('ffmpeg') and len(temp_audio_files) > 0:
#             try:
#                 # Concatenate all audio files
#                 audio_list_file = output_path.replace('.mp4', '_audio_list.txt')
#                 with open(audio_list_file, 'w', encoding='utf-8') as f:
#                     for wav_path, _ in temp_audio_files:
#                         # Use absolute path and escape for Windows
#                         abs_path = os.path.abspath(wav_path).replace('\\', '/')
#                         f.write(f"file '{abs_path}'\n")
                
#                 concat_audio_path = output_path.replace('.mp4', '_combined_audio.wav')
                
#                 logger.info("Combining audio files...")
#                 # Combine audio files
#                 result = subprocess.run([
#                     'ffmpeg', '-f', 'concat', '-safe', '0', '-i', audio_list_file,
#                     '-c', 'copy', concat_audio_path, '-y'
#                 ], capture_output=True, text=True)
                
#                 if result.returncode != 0:
#                     logger.error(f"FFmpeg audio concat failed: {result.stderr}")
#                     raise subprocess.CalledProcessError(result.returncode, 'ffmpeg')
                
#                 logger.info("Combining video with audio...")
#                 # Combine video with audio
#                 result = subprocess.run([
#                     'ffmpeg', '-i', output_path, '-i', concat_audio_path,
#                     '-c:v', 'copy', '-c:a', 'aac', '-shortest', 
#                     final_output_path, '-y'
#                 ], capture_output=True, text=True)
                
#                 if result.returncode != 0:
#                     logger.error(f"FFmpeg video+audio combine failed: {result.stderr}")
#                     raise subprocess.CalledProcessError(result.returncode, 'ffmpeg')
                
#                 # Clean up temporary files
#                 try:
#                     os.remove(output_path)  # Remove video-only file
#                     os.remove(audio_list_file)
#                     os.remove(concat_audio_path)
                    
#                     # Move final file to original name
#                     shutil.move(final_output_path, output_path)
#                 except Exception as cleanup_error:
#                     logger.warning(f"Cleanup error: {cleanup_error}")
                
#                 logger.info(f"Video with audio created successfully: {output_path}")
#                 return output_path
                
#             except subprocess.CalledProcessError as e:
#                 logger.error(f"FFmpeg error: {e}")
#                 logger.info(f"Video created without audio: {output_path}")
#                 return output_path
#             except Exception as e:
#                 logger.error(f"Unexpected error during audio-video combination: {e}")
#                 logger.info(f"Video created without audio: {output_path}")
#                 return output_path
#         else:
#             if not shutil.which('ffmpeg'):
#                 logger.warning("ffmpeg not found - video created without audio")
#             logger.info(f"Video created: {output_path}")
#             return output_path
        
#     except Exception as e:
#         logger.error(f"Failed to create video: {e}")
#         return None
#     finally:
#         # Clean up temporary audio files
#         for wav_path, original_path in temp_audio_files:
#             if wav_path != original_path and os.path.exists(wav_path):
#                 try:
#                     os.remove(wav_path)
#                 except Exception as e:
#                     logger.warning(f"Failed to clean up temp audio file {wav_path}: {e}")

# Also add these imports at the top of your file if not already present:

def get_audio_duration(audio_path):
    """Get the duration of an audio file using file size estimation."""
    try:
        if audio_path.endswith('.mp3'):
            # For MP3, estimate duration from file size
            file_size = os.path.getsize(audio_path)
            # Rough estimation: 1MB â‰ˆ 60 seconds for typical speech
            # More conservative estimate for better sync
            estimated_duration = max(3, file_size / (1024 * 15))  # Adjust divisor for accuracy
            return estimated_duration
        elif audio_path.endswith('.wav'):
            # For WAV files, try to use wave module
            import wave
            with wave.open(audio_path, 'r') as wav_file:
                frames = wav_file.getnframes()
                rate = wav_file.getframerate()
                duration = frames / float(rate)
                return duration
        else:
            return VIDEO_CONFIG["slide_duration"]  # fallback
    except Exception as e:
        logger.error(f"Failed to get audio duration: {e}")
        return VIDEO_CONFIG["slide_duration"]  # fallback duration
# === PPT TO VIDEO CONVERSION ===

# def convert_ppt_to_images(ppt_path, output_dir):
#     """Convert PowerPoint slides to images."""
#     logger.info(f"Converting PPT to images: {ppt_path}")
    
#     try:
#         # Create output directory
#         os.makedirs(output_dir, exist_ok=True)
        
#         # Use python-pptx to read the presentation
#         prs = Presentation(ppt_path)
#         image_files = []
        
#         # For each slide, we'll create a screenshot-like image
#         # Note: This is a simplified approach. For better quality, you might need
#         # to use tools like LibreOffice in headless mode or other conversion tools
        
#         for i, slide in enumerate(prs.slides):
#             image_path = os.path.join(output_dir, f"slide_{i+1:02d}.png")
            
#             # Create a blank image with presentation dimensions
#             # Convert inches to pixels (assuming 96 DPI)
#             width_px = int(prs.slide_width.inches * 96)
#             height_px = int(prs.slide_height.inches * 96)
            
#             # Scale to video resolution
#             scale_x = VIDEO_CONFIG["width"] / width_px
#             scale_y = VIDEO_CONFIG["height"] / height_px
#             scale = min(scale_x, scale_y)
            
#             final_width = int(width_px * scale)
#             final_height = int(height_px * scale)
            
#             # Create image (this is a placeholder - in practice you'd need proper PPT rendering)
#             img = Image.new('RGB', (VIDEO_CONFIG["width"], VIDEO_CONFIG["height"]), 
#                           color=VIDEO_CONFIG["background_color"])
            
#             draw = ImageDraw.Draw(img)
            
#             # Add slide number and basic layout
#             draw.rectangle([50, 50, VIDEO_CONFIG["width"]-50, VIDEO_CONFIG["height"]-50], 
#                          outline=(200, 200, 200), width=2)
            
#             # Add slide number
#             try:
#                 font = ImageDraw.ImageFont.load_default()
#                 draw.text((VIDEO_CONFIG["width"]//2, 100), f"Slide {i+1}", 
#                          fill=(50, 50, 50), font=font, anchor="mm")
#             except:
#                 draw.text((VIDEO_CONFIG["width"]//2, 100), f"Slide {i+1}", 
#                          fill=(50, 50, 50), anchor="mm")
            
#             img.save(image_path, 'PNG', quality=95)
#             image_files.append(image_path)
#             logger.info(f"Created slide image: {image_path}")
        
#         return image_files
        
#     except Exception as e:
#         logger.error(f"Failed to convert PPT to images: {e}")
#         return []
import subprocess, glob

import subprocess, glob, shutil

import subprocess, glob, os, shutil
from pdf2image import convert_from_path

def convert_ppt_to_images(ppt_path, output_dir):
    """Convert PPT to images (one PNG per slide) using LibreOffice + pdf2image."""
    logger.info(f"Converting PPT to images via PDF: {ppt_path}")
    os.makedirs(output_dir, exist_ok=True)

    try:
        # Step 1: Convert PPT â†’ PDF
        cmd = [
            "soffice", "--headless", "--convert-to", "pdf",
            "--outdir", output_dir, ppt_path
        ]
        result = subprocess.run(cmd, capture_output=True, text=True)

        if result.returncode != 0:
            logger.error(f"LibreOffice PDF export failed: {result.stderr}")
            return []

        # Step 2: Find the generated PDF
        base = os.path.splitext(os.path.basename(ppt_path))[0]
        pdf_path = os.path.join(output_dir, base + ".pdf")
        if not os.path.exists(pdf_path):
            logger.error("PDF not found after LibreOffice export")
            return []

        # Step 3: Convert PDF pages â†’ PNGs
        images = convert_from_path(pdf_path, dpi=200)
        image_files = []
        for i, page in enumerate(images, start=1):
            img_file = os.path.join(output_dir, f"slide_{i:02d}.png")
            page.save(img_file, "PNG")
            image_files.append(img_file)

        logger.info(f"Generated {len(image_files)} slide images")
        return image_files

    except Exception as e:
        logger.error(f"Failed to convert PPT to images: {e}")
        return []

# def create_video_from_slides_and_audio(image_files, audio_files, output_path):
#     """Create video from slide images and audio files."""
#     logger.info("Creating video from slides and audio...")
    
#     try:
#         video_clips = []
        
#         for i, (image_path, audio_path) in enumerate(zip(image_files, audio_files)):
#             if not os.path.exists(image_path) or not os.path.exists(audio_path):
#                 logger.warning(f"Missing file for slide {i+1}")
#                 continue
            
#             # Get audio duration
#             audio_duration = get_audio_duration(audio_path)
            
#             # Create image clip
#             img_clip = mp.ImageClip(image_path, duration=audio_duration)
            
#             # Load audio
#             audio_clip = mp.AudioFileClip(audio_path)
            
#             # Combine image and audio
#             video_clip = img_clip.set_audio(audio_clip)
            
#             # Add fade transition
#             if i > 0:
#                 video_clip = video_clip.fadein(VIDEO_CONFIG["transition_duration"])
            
#             video_clips.append(video_clip)
        
#         if not video_clips:
#             raise Exception("No valid video clips created")
        
#         # Concatenate all clips
#         final_video = mp.concatenate_videoclips(video_clips, method="compose")
        
#         # Write the video file
#         final_video.write_videofile(
#             output_path,
#             fps=VIDEO_CONFIG["fps"],
#             codec='libx264',
#             audio_codec='aac',
#             temp_audiofile='temp-audio.m4a',
#             remove_temp=True,
#             verbose=False,
#             logger=None
#         )
        
#         # Clean up
#         for clip in video_clips:
#             clip.close()
#         final_video.close()
        
#         logger.info(f"Video created successfully: {output_path}")
#         return output_path
        
#     except Exception as e:
#         logger.error(f"Failed to create video: {e}")
#         return None
def create_video_from_slides_and_audio(image_files, audio_files, output_path):
    """Create video from slide images and audio files."""
    logger.info("Creating video from slides and audio...")
    
    try:
        import moviepy.editor as mp_editor
        
        video_clips = []
        
        for i, (image_path, audio_path) in enumerate(zip(image_files, audio_files)):
            if not os.path.exists(image_path) or not os.path.exists(audio_path):
                logger.warning(f"Missing file for slide {i+1}")
                continue
            
            # Get audio duration
            audio_duration = get_audio_duration(audio_path)
            logger.info(f"Slide {i+1}: {audio_duration:.1f}s duration")
            
            # Create image clip - CORRECTED
            img_clip = mp_editor.ImageClip(image_path, duration=audio_duration)
            
            # Load audio - CORRECTED
            audio_clip = mp_editor.AudioFileClip(audio_path)
            
            # Combine image and audio - CORRECTED
            video_clip = img_clip.set_audio(audio_clip)
            
            # Add fade transition
            if i > 0:
                video_clip = video_clip.fadein(VIDEO_CONFIG["transition_duration"])
            
            video_clips.append(video_clip)
            logger.info(f"Processed slide {i+1}")
        
        if not video_clips:
            raise Exception("No valid video clips created")
        
        # Concatenate all clips
        logger.info("Concatenating video clips...")
        final_video = mp_editor.concatenate_videoclips(video_clips, method="compose")
        
        # Write the video file
        logger.info(f"Writing video to {output_path}...")
        final_video.write_videofile(
            output_path,
            fps=VIDEO_CONFIG["fps"],
            codec='libx264',
            audio_codec='aac',
            temp_audiofile='temp-audio.m4a',
            remove_temp=True,
            verbose=False,
            logger=None
        )
        
        # Clean up
        for clip in video_clips:
            clip.close()
        final_video.close()
        
        logger.info(f"Video created successfully: {output_path}")
        return output_path
        
    except ImportError:
        logger.error("MoviePy not properly installed. Please install with: pip install moviepy")
        return None
    except Exception as e:
        logger.error(f"Failed to create video: {e}")
        return None

def create_presentation_video(parsed_data, topic, ppt_path):
    """Main function to create video from presentation data."""
    logger.info("Starting video creation process...")
    
    # Create temporary directories
    temp_dir = tempfile.mkdtemp(prefix="ppt_video_")
    images_dir = os.path.join(temp_dir, "images")
    audio_dir = os.path.join(temp_dir, "audio")
    
    try:
        os.makedirs(images_dir, exist_ok=True)
        os.makedirs(audio_dir, exist_ok=True)
        
        # Step 1: Generate narration script
        narration_data = generate_narration_script(parsed_data, topic)
        
        # Step 2: Create audio files
        audio_files = []
        
        # Title slide audio
        title_audio_path = os.path.join(audio_dir, "title.mp3")
        if create_audio_from_text(narration_data["title_narration"], title_audio_path):
            audio_files.append(title_audio_path)
        
        # Content slides audio
        for i, narration in enumerate(narration_data["slide_narrations"]):
            audio_path = os.path.join(audio_dir, f"slide_{i+1:02d}.mp3")
            if create_audio_from_text(narration, audio_path):
                audio_files.append(audio_path)
        
        # Conclusion audio (if exists)
        if narration_data.get("conclusion"):
            conclusion_audio_path = os.path.join(audio_dir, "conclusion.mp3")
            if create_audio_from_text(narration_data["conclusion"], conclusion_audio_path):
                audio_files.append(conclusion_audio_path)
        
        logger.info(f"Created {len(audio_files)} audio files")
        
        # Step 3: Convert PPT to images
        image_files = convert_ppt_to_images(ppt_path, images_dir)
        
        # Ensure we have matching numbers of images and audio files
        min_length = min(len(image_files), len(audio_files))
        image_files = image_files[:min_length]
        audio_files = audio_files[:min_length]
        
        if min_length == 0:
            raise Exception("No matching image and audio files found")
        
        # Step 4: Create video
        video_filename = topic.strip().replace(" ", "_") + "_presentation.mp4"
        video_path = create_video_from_slides_and_audio(image_files, audio_files, video_filename)
        
        if video_path and os.path.exists(video_path):
            logger.info(f"Video creation successful: {video_path}")
            file_size = os.path.getsize(video_path) / (1024 * 1024)  # MB
            logger.info(f"Video file size: {file_size:.1f} MB")
            return video_path
        else:
            raise Exception("Video creation failed")
    
    except Exception as e:
        logger.error(f"Video creation process failed: {e}")
        return None
    
    finally:
        # Clean up temporary files
        try:
            shutil.rmtree(temp_dir)
            logger.info("Temporary files cleaned up")
        except:
            logger.warning("Failed to clean up temporary files")

# === UTILITY FUNCTIONS ===

def classify_paragraph_type(raw_text):
    """Classifies the type of presentation based on user input text."""
    prompt = f"""
    Classify the following text into one category:
    - business
    - academic
    - technical
    - educational
    - motivational
    - general

    Text:
    \"\"\"{raw_text}\"\"\"

    Respond with only one word (the category).
    """
    result = llm.invoke(prompt).content.strip().lower()
    if result not in ["business", "academic", "technical", "educational", "motivational", "general"]:
        return "business"  # default fallback
    return result

def refine_paragraph_input(raw_text, category="business"):
    """Refines raw text differently depending on content type."""
    style_guidelines = {
        "business": "Make it concise, factual, and strategic. Focus on insights, numbers, and implications.",
        "academic": "Make it clear, explanatory, with definitions and logical flow.",
        "technical": "Make it precise, structured, with focus on processes, systems, and technical clarity.",
        "educational": "Make it simple, clear, and beginner-friendly with learning outcomes.",
        "motivational": "Make it inspiring, story-driven, with positive tone and key messages.",
        "general": "Make it well-structured, clear, and neutral."
    }

    guidelines = style_guidelines.get(category, style_guidelines["business"])

    prompt = f"""
    Refine and summarize the following text for a professional {category} presentation.
    {guidelines}
    Output 5-6 concise paragraphs, well-structured for slide generation.
    Do not use bold, italics, underlines, Markdown symbols (like *, **, _), or any decorative formatting.
    Only plain text should be returned.
    Text:
    \"\"\"{raw_text}\"\"\"

    Refined version:
    """
    return llm.invoke(prompt).content

def get_slide_content_from_paragraph(context_text, category="business", min_slides=1, max_slides=15):
    """Generate slide content from paragraph input."""
    style_instructions = {
        "business": "Use a McKinsey-style with charts, insights, and data-driven points.",
        "academic": "Use an academic style with definitions, theories, and structured explanation.",
        "technical": "Use a technical style with system diagrams, architecture, workflows, or pseudocode.",
        "educational": "Use an educational style with simple language, learning outcomes, and key concepts.",
        "motivational": "Use a motivational style with quotes, storytelling, and call-to-action messages.",
        "general": "Use a general informative style with clarity and balanced explanation."
    }
    style = style_instructions.get(category, style_instructions["business"])

    prompt = f"""
    You are preparing a professional {category} PowerPoint presentation outline 
    based on the following refined content:

    \"\"\"{context_text}\"\"\"

    Rules:
    - Create between {min_slides} and {max_slides} slides depending on content richness.
    - {style}
    - Return ONLY valid JSON, no explanations or markdown.
    - "intro" must be maximum 2-3 lines.
    - Each slide must include:
      - "title": max 8 words
      - "insight": one or two sentences summarizing the key takeaway
      - "type": either "bullets" or "chart"
      - If "type" = "bullets", "data" must be an array of objects:
        {{"point": "short phrase", "desc": "5-6 line proper explanation"}}
      - If "type" = "chart", "data" must be JSON:
        {{
          "type": "BAR" | "LINE" | "PIE" | "COLUMN" | "DOUGHNUT" | "AREA" | "SCATTER" | "STACKED_BAR",
          "data": [["Label1", 123], ["Label2", 456]],
          "source": "Source: Organization, 2025"
        }}
        and optionally include "context": "Why this data matters"

    JSON format:
    {{
      "intro": "short intro text",
      "slides": [
        {{
          "title": "Slide Title",
          "insight": "Key insight",
          "type": "bullets",
          "data": [
            {{"point": "Clarity", "desc": "Use straightforward language to avoid confusion"}},
            {{"point": "Audience Awareness", "desc": "Adapt tone and style to the readers"}}
          ]
        }},
        {{
          "title": "Slide with Chart",
          "insight": "Data-driven point",
          "type": "chart",
          "data": {{
            "type": "COLUMN",
            "data": [["X", 10], ["Y", 20]],
            "source": "Source: Example Org, 2025"
          }},
          "context": "Why this data is important"
        }}
      ]
    }}
    """
    return llm.invoke(prompt).content

def generate_topic_from_paragraph(context_text):
    """Generate a topic title from paragraph content."""
    prompt = f"""
    Generate a concise, professional presentation title (max 6 words) from this content:
    "{context_text[:200]}..."
    
    Return only the title, no quotes or formatting.
    """
    return llm.invoke(prompt).content.strip()

# === MAIN EXECUTION ===

if __name__ == "__main__":
    print("=== PPT to Video Converter with AI Narration ===\n")
    
    mode = input("Choose input mode (1=Topic, 2=Paragraph): ").strip()
    create_video = input("Create video after PPT? (y/n): ").strip().lower() == 'y'

    if mode == "1":
        # === Topic-based workflow ===
        topic = input("Enter your presentation topic: ").strip()
        try:
            n = int(input("Number of slides: "))
            if n < 1:
                raise ValueError("Slide count must be at least 1.")
        except Exception as e:
            print(f"Invalid input: {e}")
            exit()

        print(f"\n[1/4] Generating JSON content with AI for '{topic}'...")
        raw_content = get_slide_content_with_charts(topic, n)

        print("[2/4] Parsing JSON response...")
        parsed_slides = parse_json_slides(raw_content)

        if not parsed_slides.get("slides"):
            print("Failed to generate or parse JSON content. Please try again.")
            exit()

        print(f"[3/4] Building PowerPoint presentation...")
        ppt_path = build_mckinsey_ppt(parsed_slides, topic)
        
        if create_video:
            print(f"[4/4] Creating video with AI narration...")
            video_path = create_presentation_video(parsed_slides, topic, ppt_path)
            if video_path:
                print(f"\n SUCCESS!")
                print(f" PowerPoint: {ppt_path}")
                print(f" Video: {video_path}")
            else:
                print(f"\nâš ï¸  PowerPoint created: {ppt_path}")
                print("âŒ Video creation failed")

    elif mode == "2":
        # === Paragraph-based workflow ===
        context_text = input("Enter paragraph/context for your presentation: ").strip()

        print("\n[1/6] Detecting content type...")
        category = classify_paragraph_type(context_text)
        print(f" Detected category: {category}")

        print("\n[2/6] Refining input text...")
        refined_text = refine_paragraph_input(context_text, category)

        topic = generate_topic_from_paragraph(context_text)
        print(f"â†’ Generated topic: {topic}")

        print("\n[3/6] Generating JSON slides...")
        raw_content = get_slide_content_from_paragraph(refined_text, category)

        print("[4/6] Parsing JSON response...")
        parsed_slides = parse_json_slides(raw_content)

        if not parsed_slides.get("slides"):
            print("Failed to generate or parse JSON content. Please try again.")
            exit()

        print(f"[5/6] Building PowerPoint presentation...")
        ppt_path = build_mckinsey_ppt(parsed_slides, topic)
        
        if create_video:
            print(f"[6/6] Creating video with AI narration...")
            video_path = create_presentation_video(parsed_slides, topic, ppt_path)
            if video_path:
                print(f"\n SUCCESS!")
                print(f"PowerPoint: {ppt_path}")
                print(f" Video: {video_path}")
            else:
                print(f"\nPowerPoint created: {ppt_path}")
                print(" Video creation failed")

    else:
        print("Invalid mode. Exiting.")